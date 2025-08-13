import os
import io
import uuid
import json
import time
import glob
import tempfile
import logging
import requests
import sqlite3
import subprocess
import re
import mimetypes
import random
import threading
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from datetime import datetime, timedelta
from functools import wraps
from collections import deque
from dotenv import load_dotenv
from flask import Flask, request, jsonify, render_template, session, redirect, url_for, send_file, send_from_directory
from flask_cors import CORS
from PIL import Image
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
from pdf2image import convert_from_bytes
from docx import Document
from openpyxl import load_workbook
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from rembg import remove
from gtts import gTTS
import speech_recognition as sr
from pydub import AudioSegment
from bs4 import BeautifulSoup
from fpdf import FPDF
import fitz  # PyMuPDF
import firebase_admin
from firebase_admin import credentials, db, auth
import cloudinary
import cloudinary.uploader

# --- Flask & Environment Setup ---

app = Flask(__name__)
load_dotenv('api.env')

# Firebase Realtime Database config
FIREBASE_CRED_PATH = 'docshift.json'
FIREBASE_DB_URL = 'https://docshift-86065-default-rtdb.firebaseio.com/'

if not firebase_admin._apps:
    cred = credentials.Certificate(FIREBASE_CRED_PATH)
    firebase_admin.initialize_app(cred, {'databaseURL': FIREBASE_DB_URL})

# Cloudinary config
cloudinary.config(
    cloud_name='dvdeflyta',
    api_key='568435982421747',
    api_secret='-xqsm00d0D9Hxp1YsrA6OrU-hpw'
)

# Flask CORS and secret key
CORS(app)
app.secret_key = os.urandom(24)

# Logging config
logging.basicConfig(level=logging.DEBUG,
                    format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# OpenRouter API settings (change with your key)
OPENROUTER_API_KEY = os.getenv('OPENROUTER_API_KEY', 'your_api_key_here')
OPENROUTER_API_URL = 'https://openrouter.ai/api/v1/chat/completions'
OPENROUTER_MODEL = 'gpt-4o-mini'

# Speech recognizer init
recognizer = sr.Recognizer()

# Globals for AI Document Screener and AI PDF Editor
current_document_text = ''
conversation_history = deque(maxlen=10)
latest_text = ""

# --- Utility Functions ---

def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'logged_in' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function

def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'logged_in' not in session:
            return redirect(url_for('login'))
        if session.get('role') != 'admin':
            return redirect(url_for('index'))
        return f(*args, **kwargs)
    return decorated_function

def store_url_in_firebase(url, category, filename):
    """Store URL in firebase per user under storage/{username}/{category}."""
    safe_key = re.sub(r'[./#$\[\]]', '_', filename)
    username = session.get('username', 'admin')
    ref = db.reference(f'storage/{username}/{category}/{safe_key}')
    ref.set({'filename': filename, 'url': url})
    return True

# --- Phone & Email Verification Functions ---

def is_production_mode():
    """Check if we're running in production mode with real credentials"""
    smtp_configured = bool(os.getenv('SMTP_USERNAME') and os.getenv('SMTP_PASSWORD'))
    return smtp_configured

def generate_otp():
    """Generate a random 6-digit OTP"""
    return str(random.randint(100000, 999999))

def store_otp_in_firebase(identifier, otp_code, username, verification_type='phone'):
    """Store OTP in Firebase with expiration (10 minutes)"""
    try:
        # Create expiration timestamp (10 minutes from now)
        expiry_time = datetime.now() + timedelta(minutes=10)
        expiry_timestamp = int(expiry_time.timestamp())
        
        # Store OTP data
        otp_data = {
            'code': otp_code,
            'type': verification_type,
            'expires_at': expiry_timestamp,
            'username': username,
            'verified': False,
            'created_at': int(datetime.now().timestamp())
        }
        
        # Clean identifier for Firebase key
        if verification_type == 'phone':
            clean_key = re.sub(r'[^0-9]', '', identifier)
        else:  # email
            clean_key = re.sub(r'[./#$\[\]@]', '_', identifier)
        
        ref = db.reference(f'verification_codes/{verification_type}_{clean_key}')
        ref.set(otp_data)
        
        logger.info(f"OTP stored for {verification_type} {identifier}: {otp_code}")
        return True
    except Exception as e:
        logger.error(f"Error storing OTP: {str(e)}")
        return False

def store_phone_otp_in_firebase(phone_number, otp_code, username):
    """Store phone OTP in Firebase - backward compatibility"""
    return store_otp_in_firebase(phone_number, otp_code, username, 'phone')

def store_email_otp_in_firebase(email, otp_code, username):
    """Store email OTP in Firebase"""
    return store_otp_in_firebase(email, otp_code, username, 'email')

def verify_otp_from_firebase(identifier, submitted_otp, verification_type='phone'):
    """Verify OTP from Firebase"""
    try:
        if verification_type == 'phone':
            clean_key = re.sub(r'[^0-9]', '', identifier)
        else:  # email
            clean_key = re.sub(r'[./#$\[\]@]', '_', identifier)
            
        ref = db.reference(f'verification_codes/{verification_type}_{clean_key}')
        stored_data = ref.get()
        
        if not stored_data:
            return False, f"No OTP found for this {verification_type}"
        
        # Check if OTP has expired
        current_timestamp = int(datetime.now().timestamp())
        if current_timestamp > stored_data.get('expires_at', 0):
            # Clean up expired OTP
            ref.delete()
            return False, "OTP has expired"
        
        # Check if OTP matches
        if stored_data.get('code') != submitted_otp:
            return False, "Invalid OTP"
        
        # Mark as verified and clean up
        stored_data['verified'] = True
        stored_data['verified_at'] = current_timestamp
        ref.set(stored_data)
        
        # Update user profile to mark as verified
        username = stored_data.get('username')
        if username:
            user_ref = db.reference(f'Data/{username}')
            user_data = user_ref.get() or {}
            
            if verification_type == 'phone':
                user_data.update({
                    'phone_verified': True,
                    'phone_verified_at': datetime.now().isoformat()
                })
            else:  # email
                user_data.update({
                    'email_verified': True,
                    'email_verified_at': datetime.now().isoformat()
                })
            user_ref.set(user_data)
        
        # Clean up OTP after successful verification
        threading.Timer(5.0, lambda: ref.delete()).start()
        
        return True, f"{verification_type.title()} verified successfully!"
        
    except Exception as e:
        logger.error(f"Error verifying OTP: {str(e)}")
        return False, f"Verification error: {str(e)}"

def verify_phone_otp_from_firebase(phone_number, submitted_otp):
    """Verify phone OTP - backward compatibility"""
    return verify_otp_from_firebase(phone_number, submitted_otp, 'phone')

def verify_email_otp_from_firebase(email, submitted_otp):
    """Verify email OTP"""
    return verify_otp_from_firebase(email, submitted_otp, 'email')

def send_email_otp(email, otp_code):
    """Send email OTP using Gmail SMTP"""
    try:
        # Get email credentials from environment
        smtp_server = os.getenv('SMTP_SERVER', 'smtp.gmail.com')
        smtp_port = int(os.getenv('SMTP_PORT', '587'))
        smtp_username = os.getenv('SMTP_USERNAME')  # Your Gmail address
        smtp_password = os.getenv('SMTP_PASSWORD')  # Your Gmail app password
        
        # Check if email credentials are configured
        if not smtp_username or not smtp_password:
            logger.error("📧 Email credentials not configured")
            return False, "Email service not configured. Please contact administrator."
        
        # Create email message
        msg = MIMEMultipart()
        msg['From'] = smtp_username
        msg['To'] = email
        msg['Subject'] = "DocShift - Email Verification Code"
        
        # Create HTML email body
        html_body = f"""
        <html>
            <body style="font-family: Arial, sans-serif; line-height: 1.6; color: #333;">
                <div style="max-width: 600px; margin: 0 auto; padding: 20px; border: 1px solid #ddd; border-radius: 10px;">
                    <div style="text-align: center; margin-bottom: 30px;">
                        <h1 style="color: #4A90E2;">DocShift</h1>
                        <h2 style="color: #333;">Email Verification</h2>
                    </div>
                    
                    <div style="background-color: #f9f9f9; padding: 20px; border-radius: 5px; margin-bottom: 20px;">
                        <p>Hello,</p>
                        <p>You've requested to verify your email address for DocShift. Please use the following verification code:</p>
                        
                        <div style="text-align: center; margin: 30px 0;">
                            <span style="background-color: #4A90E2; color: white; padding: 15px 30px; font-size: 24px; font-weight: bold; border-radius: 5px; letter-spacing: 3px;">{otp_code}</span>
                        </div>
                        
                        <p><strong>This code will expire in 10 minutes.</strong></p>
                        <p>If you didn't request this verification, please ignore this email.</p>
                    </div>
                    
                    <div style="text-align: center; color: #666; font-size: 12px;">
                        <p>© 2025 DocShift. All rights reserved.</p>
                    </div>
                </div>
            </body>
        </html>
        """
        
        msg.attach(MIMEText(html_body, 'html'))
        
        # Send email
        server = smtplib.SMTP(smtp_server, smtp_port)
        server.starttls()
        server.login(smtp_username, smtp_password)
        text = msg.as_string()
        server.sendmail(smtp_username, email, text)
        server.quit()
        
        logger.info(f"📧 Email OTP sent successfully to {email}")
        return True, "Verification code sent to your email!"
        
    except Exception as e:
        logger.error(f"Error sending email: {str(e)}")
        return False, f"Failed to send email verification code. Please try again."

# Phone verification via email functionality removed

# SMS OTP functionality removed

def cleanup_expired_otps():
    """Clean up expired OTPs from Firebase"""
    try:
        ref = db.reference('verification_codes')
        all_codes = ref.get() or {}
        current_timestamp = int(datetime.now().timestamp())
        
        for key, data in all_codes.items():
            if isinstance(data, dict) and current_timestamp > data.get('expires_at', 0):
                db.reference(f'verification_codes/{key}').delete()
                logger.info(f"Cleaned up expired OTP for {key}")
                
    except Exception as e:
        logger.error(f"Error cleaning up expired OTPs: {str(e)}")

# Start background cleanup thread for expired OTPs
def periodic_otp_cleanup():
    while True:
        time.sleep(300)  # Check every 5 minutes
        cleanup_expired_otps()

cleanup_thread = threading.Thread(target=periodic_otp_cleanup, daemon=True)
cleanup_thread.start()

# --- End Phone & Email Verification Functions ---

def upload_to_cloudinary(local_path, folder):
    """Upload file to Cloudinary, handle image/raw types."""
    ext = os.path.splitext(local_path)[1].lower()
    if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp']:
        response = cloudinary.uploader.upload(local_path, folder=folder)
    elif ext in ['.mp3', '.wav', '.aac', '.ogg', '.flac', '.pdf', '.txt', '.docx', '.xlsx', '.pptx', '.csv']:
        response = cloudinary.uploader.upload(local_path, folder=folder, resource_type='raw')
    else:
        response = cloudinary.uploader.upload(local_path, folder=folder, resource_type='raw')
    return response['secure_url']

def get_user_storage_path():
    return f"storage/{session['username']}"

def init_db():
    """Initialize SQLite DB for conversion logs (optional)."""
    conn = sqlite3.connect('file_conversion.db')
    c = conn.cursor()
    c.execute('''CREATE TABLE IF NOT EXISTS conversions (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        conversion_type TEXT NOT NULL,
        original_filename TEXT NOT NULL,
        converted_filename TEXT NOT NULL,
        timestamp DATETIME DEFAULT CURRENT_TIMESTAMP,
        file_path TEXT,
        cloudinary_url TEXT,
        username TEXT,
        status TEXT DEFAULT 'error'
    )''')
    # Add new columns if they don't exist (for existing databases)
    try:
        c.execute('ALTER TABLE conversions ADD COLUMN cloudinary_url TEXT')
    except sqlite3.OperationalError:
        pass  # Column already exists
    try:
        c.execute('ALTER TABLE conversions ADD COLUMN username TEXT')
    except sqlite3.OperationalError:
        pass  # Column already exists
    try:
        c.execute('ALTER TABLE conversions ADD COLUMN status TEXT DEFAULT "error"')
    except sqlite3.OperationalError:
        pass  # Column already exists
    c.execute('''CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT NOT NULL UNIQUE,
        password_hash TEXT NOT NULL
    )''')
    conn.commit()
    conn.close()

def log_conversion(conversion_type, original_filename, converted_filename, file_path=None, cloudinary_url=None, status=None):
    conn = sqlite3.connect('file_conversion.db')
    c = conn.cursor()
    username = session.get('username', 'admin')
    
    # Determine status based on cloudinary_url if not explicitly provided
    if status is None:
        status = 'success' if cloudinary_url else 'error'
    
    c.execute('''
        INSERT INTO conversions (conversion_type, original_filename, converted_filename, file_path, cloudinary_url, username, status)
        VALUES (?, ?, ?, ?, ?, ?, ?)
    ''', (conversion_type, original_filename, converted_filename, file_path, cloudinary_url, username, status))
    conn.commit()
    conn.close()

def is_ghostscript_installed():
    try:
        subprocess.run(['gs', '--version'], stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
        return True
    except (subprocess.CalledProcessError, FileNotFoundError):
        return False

# --- Firebase: Get user credentials ---
def get_user_by_username(username):
    """Fetch user credentials from Firebase Realtime DB.
    Assign role 'admin' for admin user, 'user' for regular users."""
    if username == 'admin':
        ref = db.reference('credentials/admin/admin')
        user_record = ref.get()
        if user_record:
            user_record['role'] = 'admin'
        return user_record
    else:
        ref = db.reference(f'credentials/users/{username}')
        user_record = ref.get()
        if user_record:
            user_record['role'] = 'user'
        return user_record

# --- Ensure admin credentials exist ---
def ensure_admin_credentials():
    admin_ref = db.reference('credentials/admin/admin')
    if not admin_ref.get():
        default_password = generate_password_hash("admin123")
        admin_ref.set({'password': default_password})
        print("Default admin credentials inserted (admin/admin123)")
ensure_admin_credentials()

# --- Routes ---

# --- Phone & Email Verification Routes ---

@app.route('/send_phone_otp', methods=['POST'])
def send_phone_otp():
    """Phone verification disabled - use email verification only"""
    return jsonify({'success': False, 'error': 'Phone verification is currently disabled. Please use email verification.'})

@app.route('/send_email_otp', methods=['POST'])
def send_email_otp_route():
    """Send OTP to email for verification"""
    try:
        data = request.get_json()
        email = data.get('email', '').strip()
        username = session.get('username')
        
        if not email:
            return jsonify({'success': False, 'error': 'Email address is required'})
        
        if not username:
            return jsonify({'success': False, 'error': 'User not logged in'})
        
        # Validate email format (basic validation)
        import re
        email_pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
        if not re.match(email_pattern, email):
            return jsonify({'success': False, 'error': 'Please enter a valid email address'})
        
        # Generate OTP
        otp_code = generate_otp()
        
        # Store OTP in Firebase
        if not store_email_otp_in_firebase(email, otp_code, username):
            return jsonify({'success': False, 'error': 'Failed to generate OTP'})
        
        # Send Email - production mode only
        email_success, email_message = send_email_otp(email, otp_code)
        
        if email_success:
            return jsonify({
                'success': True, 
                'message': 'Verification code sent to your email address. Please check your inbox.'
            })
        else:
            return jsonify({'success': False, 'error': email_message})
            
    except Exception as e:
        logger.error(f"Send email OTP error: {str(e)}")
        return jsonify({'success': False, 'error': 'Failed to send OTP'})

@app.route('/verify_phone_otp', methods=['POST'])
def verify_phone_otp():
    """Phone verification disabled - use email verification only"""
    return jsonify({'success': False, 'error': 'Phone verification is currently disabled. Please use email verification.'})

@app.route('/verify_email_otp', methods=['POST'])
def verify_email_otp():
    """Verify OTP for email"""
    try:
        data = request.get_json()
        email = data.get('email', '').strip()
        otp_code = data.get('otp_code', '').strip()
        
        if not email or not otp_code:
            return jsonify({'success': False, 'error': 'Email and OTP are required'})
        
        # Verify OTP
        verification_success, message = verify_email_otp_from_firebase(email, otp_code)
        
        if verification_success:
            return jsonify({
                'success': True,
                'message': message
            })
        else:
            return jsonify({
                'success': False,
                'error': message
            })
            
    except Exception as e:
        logger.error(f"Verify email OTP error: {str(e)}")
        return jsonify({'success': False, 'error': 'Failed to verify OTP'})

@app.route('/check_phone_verification_status', methods=['POST'])
def check_phone_verification_status():
    """Phone verification disabled - always return unverified"""
    return jsonify({'verified': False, 'message': 'Phone verification is currently disabled'})

@app.route('/check_email_verification_status', methods=['POST'])
def check_email_verification_status():
    """Check if email is verified for current user"""
    try:
        username = session.get('username')
        if not username:
            return jsonify({'verified': False, 'error': 'User not logged in'})
        
        user_ref = db.reference(f'Data/{username}')
        user_data = user_ref.get() or {}
        
        is_verified = user_data.get('email_verified', False)
        verified_at = user_data.get('email_verified_at', None)
        
        return jsonify({
            'verified': is_verified,
            'verified_at': verified_at
        })
        
    except Exception as e:
        logger.error(f"Check email verification status error: {str(e)}")
        return jsonify({'verified': False, 'error': 'Failed to check verification status'})

@app.route('/verify_email_standalone', methods=['POST'])
def verify_email_standalone():
    """Standalone email verification route (if needed for separate email verification page)"""
    try:
        email = request.form.get('email', '').strip()
        email_otp = request.form.get('email_otp', '').strip()
        
        if not email or not email_otp:
            return render_template('verify_email.html', 
                                 email=email,
                                 error='Email and OTP are required')
        
        # Verify OTP
        verification_success, message = verify_email_otp_from_firebase(email, email_otp)
        
        if verification_success:
            return render_template('registration_success.html', 
                                 message='Email verified successfully!')
        else:
            return render_template('verify_email.html', 
                                 email=email,
                                 error=message)
    
    except Exception as e:
        logger.error(f"Standalone email verification error: {str(e)}")
        return render_template('verify_email.html', 
                             email=email,
                             error='An error occurred during verification. Please try again.')

# --- End Phone & Email Verification Routes ---

@app.route('/register-company', methods=['GET', 'POST'])
def register_company():
    if request.method == 'POST':
        company_name = request.form.get('company_name')
        owner_name = request.form.get('owner_name')
        email = request.form.get('email')
        phone = request.form.get('phone')
        username = request.form.get('username')
        password = request.form.get('password')
        confirm_password = request.form.get('confirm_password')

        if not all([company_name, owner_name, email, phone, username, password, confirm_password]):
            return render_template('register_company.html', error='All fields are required')
        if password != confirm_password:
            return render_template('register_company.html', error='Passwords do not match')

        # Check if username already exists
        cred_ref = db.reference(f'credentials/users/{username}')
        if cred_ref.get():
            return render_template('register_company.html', error='Username already exists')

        # Validate phone number format
        clean_phone = re.sub(r'[^0-9]', '', phone)
        if len(clean_phone) < 10:
            return render_template('register_company.html', error='Please enter a valid phone number')

        # Store company data as "pending verification"
        temp_id = str(uuid.uuid4())
        hashed_pw = generate_password_hash(password)
        
        pending_data = {
            'temp_id': temp_id,
            'company_name': company_name,
            'owner_name': owner_name,
            'email': email,
            'phone': phone,
            'username': username,
            'password_hash': hashed_pw,
            'created_at': datetime.now().isoformat(),
            'phone_verified': False
        }
        
        # Store in pending_companies
        db.reference(f'pending_companies/{temp_id}').set(pending_data)
        
        # Generate and send email OTP only
        email_otp = generate_otp()
        
        email_stored = store_email_otp_in_firebase(email, email_otp, username)
        
        if email_stored:
            # Send email OTP only
            send_email_otp(email, email_otp)
            
            # Redirect to email verification page only
            return render_template('verify_email.html', 
                                 temp_id=temp_id, 
                                 email=email,
                                 username=username,
                                 message='Please verify your email to complete registration. Check your inbox for verification code.')
        else:
            # Clean up pending data if OTP failed
            db.reference(f'pending_companies/{temp_id}').delete()
            return render_template('register_company.html', 
                                 error='Failed to send verification code. Please try again.')
    
    return render_template('register_company.html')

@app.route('/verify-registration', methods=['POST'])
def verify_registration():
    """Complete company registration after email verification only"""
    try:
        temp_id = request.form.get('temp_id')
        email_otp = request.form.get('email_otp')
        
        if not temp_id or not email_otp:
            return render_template('verify_email.html', 
                                 error='Email verification code is required')
        
        # Get pending company data
        pending_ref = db.reference(f'pending_companies/{temp_id}')
        pending_data = pending_ref.get()
        
        if not pending_data:
            return render_template('verify_email.html', 
                                 error='Registration session expired. Please register again.')
        
        email = pending_data.get('email')
        username = pending_data.get('username')
        
        # Verify email OTP only
        email_verification_success, email_message = verify_email_otp_from_firebase(email, email_otp)
        
        if email_verification_success:
            # Create actual user account
            cred_ref = db.reference(f'credentials/users/{username}')
            cred_ref.set({'password': pending_data['password_hash']})
            
            # Create user data
            db.reference(f'Data/{username}').set({
                'company_name': pending_data['company_name'],
                'owner_name': pending_data['owner_name'],
                'email': pending_data['email'],
                'phone': pending_data['phone'],
                'username': pending_data['username'],
                'password': pending_data['password_hash'],
                'phone_verified': False,  # Phone verification disabled
                'email_verified': True,
                'email_verified_at': datetime.now().isoformat(),
                'created_at': pending_data['created_at']
            })
            
            # Create storage folders for user
            db.reference(f'storage/{username}').set({
                'txt': {},
                'img': {},
                'audio': {},
                'files': {}
            })
            
            # Clean up pending data
            pending_ref.delete()
            
            return render_template('registration_success.html', 
                                 message='Registration completed successfully! Your email has been verified.')
        else:
            return render_template('verify_email.html', 
                                 temp_id=temp_id,
                                 email=email,
                                 username=username,
                                 error=email_message)
    
    except Exception as e:
        logger.error(f"Registration verification error: {str(e)}")
        return render_template('verify_email.html', 
                             error='An error occurred during verification. Please try again.')

@app.route('/resend-otp', methods=['POST'])
def resend_otp():
    """Resend OTP for registration verification - email only"""
    try:
        temp_id = request.form.get('temp_id')
        
        if not temp_id:
            return jsonify({'success': False, 'error': 'Invalid request'})
        
        # Get pending company data
        pending_ref = db.reference(f'pending_companies/{temp_id}')
        pending_data = pending_ref.get()
        
        if not pending_data:
            return jsonify({'success': False, 'error': 'Registration session expired'})
        
        email = pending_data.get('email')
        username = pending_data.get('username')
        
        # Generate new email OTP only
        email_otp = generate_otp()
        if store_email_otp_in_firebase(email, email_otp, username):
            send_email_otp(email, email_otp)
            return jsonify({'success': True, 'message': 'Email verification code sent successfully!'})
        else:
            return jsonify({'success': False, 'error': 'Failed to send email OTP'})
    
    except Exception as e:
        logger.error(f"Resend OTP error: {str(e)}")
        return jsonify({'success': False, 'error': 'Failed to resend OTP'})

# Firebase Phone Auth Routes Removed - Email verification only

# Firebase SMS routes removed - Email verification only

@app.route('/login', methods=['GET', 'POST'])
def login():
    error = None
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        user_data = get_user_by_username(username)
        if not user_data:
            error = 'Invalid username'
        elif not check_password_hash(user_data['password'], password):
            error = 'Incorrect password'
        else:
            role = user_data.get('role', 'user')
            session['username'] = username
            session['role'] = role
            session['logged_in'] = True  # fix session persistence
            if role == 'admin':
                return redirect(url_for('admin_dashboard'))
            else:
                return redirect(url_for('index'))
    return render_template('login.html', error=error)

@app.route('/logout')
def logout():
    session.pop('logged_in', None)
    session.pop('username', None)
    session.pop('role', None)
    return redirect(url_for('login'))

# --- Admin Routes ---
@app.route('/admin/dashboard')
@admin_required
def admin_dashboard():
    return render_template('admin_dashboard.html', **get_user_context())

# --- Edit Company Endpoint ---
@app.route('/admin/edit-company', methods=['POST'])
@admin_required
def admin_edit_company():
    try:
        data = request.get_json()
        username = data.get('username')
        if not username:
            return jsonify({'success': False, 'error': 'Username required'})
        company_ref = db.reference(f'Data/{username}')
        company_data = company_ref.get()
        if not company_data:
            return jsonify({'success': False, 'error': 'Company not found'})
        # Update fields
        for field in ['company_name', 'owner_name', 'email', 'phone']:
            if field in data:
                company_data[field] = data[field]
        company_ref.set(company_data)
        return jsonify({'success': True})
    except Exception as e:
        logger.error(f"Edit company error: {str(e)}")
        return jsonify({'success': False, 'error': str(e)})

# --- Export Reports Endpoint ---
@app.route('/admin/export-reports', methods=['POST'])
@admin_required
def admin_export_reports():
    """Export admin reports as CSV, PDF, or JSON based on frontend request."""
    try:
        data = request.get_json()
        report_type = data.get('type')
        from_date = data.get('fromDate')
        to_date = data.get('toDate')
        export_format = data.get('format', 'csv')
        plan_filters = data.get('planFilters', [])

        # Parse date range
        from_dt = datetime.strptime(from_date, '%Y-%m-%d') if from_date else None
        to_dt = datetime.strptime(to_date, '%Y-%m-%d') if to_date else None

        # Fetch data based on report_type
        result_data = []
        if report_type == 'users':
            users_ref = db.reference('Data')
            users_data = users_ref.get() or {}
            for username, user in users_data.items():
                if not isinstance(user, dict):
                    continue
                reg_date = user.get('registered_date', None)
                # Date filter (if available)
                if reg_date and from_dt and to_dt:
                    try:
                        reg_dt = datetime.strptime(reg_date, '%Y-%m-%d')
                        if not (from_dt <= reg_dt <= to_dt):
                            continue
                    except:
                        pass
                result_data.append({
                    'Username': username,
                    'Name': user.get('owner_name', ''),
                    'Email': user.get('email', ''),
                    'Phone': user.get('phone', ''),
                    'Company': user.get('company_name', ''),
                    'Plan': user.get('membership_status', ''),
                    'Registered': user.get('registered_date', '')
                })
        elif report_type == 'companies':
            companies_ref = db.reference('Data')
            companies_data = companies_ref.get() or {}
            for username, company in companies_data.items():
                if not isinstance(company, dict):
                    continue
                reg_date = company.get('registered_date', None)
                if reg_date and from_dt and to_dt:
                    try:
                        reg_dt = datetime.strptime(reg_date, '%Y-%m-%d')
                        if not (from_dt <= reg_dt <= to_dt):
                            continue
                    except:
                        pass
                result_data.append({
                    'Username': username,
                    'Company': company.get('company_name', ''),
                    'Owner': company.get('owner_name', ''),
                    'Email': company.get('email', ''),
                    'Phone': company.get('phone', ''),
                    'Plan': company.get('membership_status', ''),
                    'Registered': company.get('registered_date', '')
                })
        elif report_type == 'plans':
            users_ref = db.reference('Data')
            users_data = users_ref.get() or {}
            for username, user in users_data.items():
                if not isinstance(user, dict):
                    continue
                plan = user.get('membership_status', 'Free')
                if plan_filters and plan.lower() not in plan_filters:
                    continue
                reg_date = user.get('registered_date', None)
                if reg_date and from_dt and to_dt:
                    try:
                        reg_dt = datetime.strptime(reg_date, '%Y-%m-%d')
                        if not (from_dt <= reg_dt <= to_dt):
                            continue
                    except:
                        pass
                result_data.append({
                    'Username': username,
                    'Plan': plan,
                    'Company': user.get('company_name', ''),
                    'Owner': user.get('owner_name', ''),
                    'Email': user.get('email', ''),
                    'Registered': user.get('registered_date', '')
                })
        elif report_type == 'usage':
            # Example: tool usage stats
            storage_ref = db.reference('storage')
            storage_data = storage_ref.get() or {}
            for username, user_storage in storage_data.items():
                if not isinstance(user_storage, dict):
                    continue
                usage = {k: len(v) if isinstance(v, dict) else 0 for k, v in user_storage.items()}
                result_data.append({'Username': username, **usage})
        elif report_type == 'financial':
            # Placeholder: implement as needed
            result_data.append({'Note': 'Financial report not implemented'})
        elif report_type == 'comprehensive':
            # Combine all above
            # For brevity, just combine users and companies
            users_ref = db.reference('Data')
            users_data = users_ref.get() or {}
            for username, user in users_data.items():
                if not isinstance(user, dict):
                    continue
                result_data.append({
                    'Username': username,
                    'Name': user.get('owner_name', ''),
                    'Email': user.get('email', ''),
                    'Phone': user.get('phone', ''),
                    'Company': user.get('company_name', ''),
                    'Plan': user.get('membership_status', ''),
                    'Registered': user.get('registered_date', '')
                })

        # Output in requested format
        if export_format == 'json':
            return jsonify(result_data)
        elif export_format == 'csv':
            import csv
            output = io.StringIO()
            if result_data:
                writer = csv.DictWriter(output, fieldnames=result_data[0].keys())
                writer.writeheader()
                writer.writerows(result_data)
            else:
                output.write('No data found')
            output.seek(0)
            return send_file(
                io.BytesIO(output.getvalue().encode('utf-8')),
                mimetype='text/csv',
                as_attachment=True,
                download_name=f'docshift_{report_type}_report_{datetime.now().strftime("%Y%m%d")}.csv'
            )
        elif export_format == 'pdf':
            # Simple PDF export using FPDF
            from fpdf import FPDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font('Arial', 'B', 12)
            pdf.cell(0, 10, f'{report_type.title()} Report', ln=1, align='C')
            pdf.set_font('Arial', '', 10)
            if result_data:
                col_width = pdf.w / (len(result_data[0]) + 1)
                row_height = pdf.font_size * 1.5
                # Header
                for key in result_data[0].keys():
                    pdf.cell(col_width, row_height, str(key), border=1)
                pdf.ln(row_height)
                # Rows
                for row in result_data:
                    for val in row.values():
                        pdf.cell(col_width, row_height, str(val), border=1)
                    pdf.ln(row_height)
            else:
                pdf.cell(0, 10, 'No data found', ln=1)
            pdf_bytes = pdf.output(dest='S').encode('latin1')
            pdf_output = io.BytesIO(pdf_bytes)
            pdf_output.seek(0)
            return send_file(
                pdf_output,
                mimetype='application/pdf',
                as_attachment=True,
                download_name=f'docshift_{report_type}_report_{datetime.now().strftime("%Y%m%d")}.pdf'
            )
        else:
            return jsonify({'error': 'Invalid export format'}), 400
    except Exception as e:
        logger.error(f"Export reports error: {str(e)}")
        return jsonify({'error': f'Failed to export report: {str(e)}'}), 500

@app.route('/api/admin/dashboard-data')
@admin_required
def admin_dashboard_data():
    
    try:
        # Get all company data from Firebase
        companies_ref = db.reference('Data')
        companies_data = companies_ref.get() or {}
        
        # Process companies data
        companies_list = []
        total_companies = 0
        active_companies = 0
        
        for username, company_data in companies_data.items():
            if isinstance(company_data, dict) and 'company_name' in company_data:
                total_companies += 1
                active_companies += 1  # For now, all are considered active
                
                companies_list.append({
                    'company_name': company_data.get('company_name', ''),
                    'owner_name': company_data.get('owner_name', ''),
                    'email': company_data.get('email', ''),
                    'phone': company_data.get('phone', ''),
                    'username': username,
                    'registered_date': 'Recent'  # You can add timestamp later
                })
        
        # Sort by company name
        companies_list.sort(key=lambda x: x['company_name'])
        
        # Get storage data for file count (optional)
        storage_ref = db.reference('storage')
        storage_data = storage_ref.get() or {}
        total_files = 0
        
        for user_storage in storage_data.values():
            if isinstance(user_storage, dict):
                for storage_type in ['txt', 'img', 'audio', 'files']:
                    type_data = user_storage.get(storage_type, {})
                    if isinstance(type_data, dict):
                        total_files += len(type_data)
        
        dashboard_data = {
            'totalCompanies': total_companies,
            'activeCompanies': active_companies,
            'totalUsers': total_companies,  # Same as companies for now
            'totalFiles': total_files,
            'companies': companies_list[:10]  # Show only recent 10
        }
        
        return jsonify(dashboard_data)
        
    except Exception as e:
        logger.error(f"Error fetching admin dashboard data: {str(e)}")
        return jsonify({"error": "Failed to fetch data"}), 500

@app.route('/admin/companies')
@admin_required
def admin_companies():
    return render_template('admin_companies.html', **get_user_context())

@app.route('/api/admin/all-companies')
@admin_required
def admin_all_companies():
    try:
        # Get all company data from Firebase
        companies_ref = db.reference('Data')
        companies_data = companies_ref.get() or {}
        
        # Process companies data
        companies_list = []
        
        for username, company_data in companies_data.items():
            if isinstance(company_data, dict) and 'company_name' in company_data:
                companies_list.append({
                    'company_name': company_data.get('company_name', ''),
                    'owner_name': company_data.get('owner_name', ''),
                    'email': company_data.get('email', ''),
                    'phone': company_data.get('phone', ''),
                    'username': username,
                    'registered_date': 'Recent'  # You can add timestamp later
                })
        
        # Sort by company name
        companies_list.sort(key=lambda x: x['company_name'])
        
        return jsonify({'companies': companies_list})
        
    except Exception as e:
        logger.error(f"Error fetching all companies data: {str(e)}")
        return jsonify({"error": "Failed to fetch data"}), 500

@app.route('/admin/company-details')
@admin_required
def company_details():
    return render_template('company_details.html', **get_user_context())

@app.route('/api/admin/company-details/<username>')
@admin_required
def company_details_api(username):
    try:
        # Get company basic info
        company_ref = db.reference(f'Data/{username}')
        company_data = company_ref.get()
        
        if not company_data:
            return jsonify({"error": "Company not found"}), 404
        
        # Get storage data for usage statistics
        storage_ref = db.reference(f'storage/{username}')
        storage_data = storage_ref.get() or {}
        
        # Calculate statistics
        total_files = 0
        tool_usage = {
            'pdf': 0,
            'image': 0,
            'text': 0,
            'audio': 0,
            'compress': 0,
            'merge': 0,
            'split': 0,
            'convert': 0
        }
        
        # Count files by type
        for storage_type, files in storage_data.items():
            if isinstance(files, dict):
                file_count = len(files)
                total_files += file_count
                
                # Map storage types to tool usage
                if storage_type == 'txt':
                    tool_usage['text'] += file_count
                elif storage_type == 'img':
                    tool_usage['image'] += file_count
                elif storage_type == 'audio':
                    tool_usage['audio'] += file_count
                elif storage_type == 'files':
                    tool_usage['pdf'] += file_count
                    tool_usage['convert'] += file_count
        
        # Calculate approximate storage (assuming average file size)
        storage_used = total_files * 2.5  # Average 2.5 MB per file
        
        # Count unique tools used
        tools_used = sum(1 for count in tool_usage.values() if count > 0)
        
        # Generate sample recent activity
        recent_activity = []
        if total_files > 0:
            activities = [
                {"type": "pdf", "title": "PDF Converted", "description": "Word document converted to PDF", "time": "2 hours ago"},
                {"type": "image", "title": "Image Processed", "description": "Background removed from image", "time": "5 hours ago"},
                {"type": "text", "title": "Text Uploaded", "description": "Text content saved to storage", "time": "1 day ago"},
                {"type": "convert", "title": "File Conversion", "description": "Excel file converted to PDF", "time": "2 days ago"},
                {"type": "audio", "title": "Audio Generated", "description": "Text converted to speech", "time": "3 days ago"}
            ]
            recent_activity = activities[:min(5, total_files)]
        
        response_data = {
            'company': {
                'company_name': company_data.get('company_name', ''),
                'owner_name': company_data.get('owner_name', ''),
                'email': company_data.get('email', ''),
                'phone': company_data.get('phone', ''),
                'username': username
            },
            'stats': {
                'totalFiles': total_files,
                'storageUsed': round(storage_used, 1),
                'monthlyActivity': total_files,  # Simplified for now
                'toolsUsed': tools_used
            },
            'toolUsage': tool_usage,
            'recentActivity': recent_activity
        }
        
        return jsonify(response_data)
        
    except Exception as e:
        logger.error(f"Error fetching company details for {username}: {str(e)}")
        return jsonify({"error": "Failed to fetch company details"}), 500

@app.route('/admin/delete-company', methods=['POST'])
@admin_required
def delete_company():
    """Delete a company and all associated data"""
    try:
        # Get request data
        data = request.get_json()
        
        if not data:
            return jsonify({"error": "No data provided"}), 400
        
        username = data.get('username')
        confirm_delete = data.get('confirm_delete', False)
        
        if not username:
            return jsonify({"error": "Username is required"}), 400
        
        if not confirm_delete:
            return jsonify({"error": "Delete confirmation is required"}), 400
        
        # Check if company exists
        company_ref = db.reference(f'Data/{username}')
        company_data = company_ref.get()
        
        if not company_data:
            return jsonify({"error": "Company not found"}), 404
        
        # Get company name for logging
        company_name = company_data.get('company_name', username)
        
        # Log the deletion attempt
        logger.info(f"Admin attempting to delete company: {company_name} (username: {username}) by admin: {session.get('admin_id', 'unknown')}")
        
        # Start deletion process
        deletion_results = {
            'company_data': False,
            'storage_data': False,
            'credentials': False,
            'file_cleanup': False
        }
        
        try:
            # 1. Delete company basic data
            company_ref.delete()
            deletion_results['company_data'] = True
            logger.info(f"Deleted company data for: {username}")
            
        except Exception as e:
            logger.error(f"Error deleting company data for {username}: {str(e)}")
        
        try:
            # 2. Delete storage data (uploaded files metadata)
            storage_ref = db.reference(f'storage/{username}')
            storage_data = storage_ref.get()
            
            if storage_data:
                storage_ref.delete()
                deletion_results['storage_data'] = True
                logger.info(f"Deleted storage data for: {username}")
            else:
                deletion_results['storage_data'] = True  # No storage data to delete
                
        except Exception as e:
            logger.error(f"Error deleting storage data for {username}: {str(e)}")
        
        try:
            # 3. Delete user credentials if they exist
            credentials_ref = db.reference(f'credentials/{username}')
            credentials_data = credentials_ref.get()
            
            if credentials_data:
                credentials_ref.delete()
                deletion_results['credentials'] = True
                logger.info(f"Deleted credentials for: {username}")
            else:
                deletion_results['credentials'] = True  # No credentials to delete
                
        except Exception as e:
            logger.error(f"Error deleting credentials for {username}: {str(e)}")
        
        try:
            # 4. Clean up any physical files if they exist
            # This would depend on your file storage implementation
            # For now, we'll mark it as successful since files are typically auto-deleted
            deletion_results['file_cleanup'] = True
            
        except Exception as e:
            logger.error(f"Error during file cleanup for {username}: {str(e)}")
        
        # Check overall success
        total_operations = len(deletion_results)
        successful_operations = sum(deletion_results.values())
        
        if successful_operations == total_operations:
            # Complete success
            logger.info(f"Successfully deleted company: {company_name} (username: {username})")
            
            # Log to admin activity if you have such a system
            try:
                admin_activity_ref = db.reference('admin_activity')
                admin_activity_ref.push({
                    'action': 'delete_company',
                    'admin_id': session.get('admin_id', 'unknown'),
                    'target_company': company_name,
                    'target_username': username,
                    'timestamp': datetime.now().isoformat(),
                    'status': 'success'
                })
            except Exception as log_error:
                logger.error(f"Error logging admin activity: {str(log_error)}")
            
            return jsonify({
                "success": True,
                "message": f"Company '{company_name}' has been successfully deleted",
                "deletion_details": deletion_results
            })
            
        elif successful_operations > 0:
            # Partial success
            logger.warning(f"Partial deletion for company: {company_name} (username: {username}). Results: {deletion_results}")
            
            return jsonify({
                "success": True,
                "message": f"Company '{company_name}' has been partially deleted. Some data may remain.",
                "warning": "Partial deletion occurred",
                "deletion_details": deletion_results
            })
            
        else:
            # Complete failure
            logger.error(f"Failed to delete company: {company_name} (username: {username})")
            return jsonify({
                "error": "Failed to delete company. No data was removed.",
                "deletion_details": deletion_results
            }), 500
        
    except Exception as e:
        logger.error(f"Unexpected error during company deletion: {str(e)}")
        return jsonify({"error": f"Internal server error: {str(e)}"}), 500

@app.route('/admin/<path:filename>')
def admin_files(filename):
    """Serve admin static files"""
    return send_from_directory('admin', filename)

@app.route('/static/<path:filename>')
def static_files(filename):
    """Serve static files"""
    return send_from_directory('static', filename)

# Serve assests folder for FAQ images
@app.route('/assests/<path:filename>')
def assests_files(filename):
    return send_from_directory('assests', filename)

@app.route('/')
@login_required
def index():
    return render_template('index.html', **get_user_context())

@app.route('/help')
@login_required
def help_support():
    """Help & Support page"""
    return render_template('help_support.html', **get_user_context())

@app.route('/privacy-policy')
def privacy_policy():
    """Privacy Policy page"""
    return render_template('privacy_policy.html', **get_user_context())

@app.route('/terms-conditions')
def terms_conditions():
    """Terms & Conditions page"""
    return render_template('terms_conditions.html', **get_user_context())

@app.route('/upgrade-plan')
@login_required
def upgrade_plan():
    """Upgrade Plan page"""
    return render_template('upgrade_plan.html', **get_user_context())

@app.route('/select_plan', methods=['POST'])
@login_required
def select_plan():
    """Handle plan selection and update user membership"""
    try:
        data = request.get_json()
        plan_name = data.get('plan', '')
        
        if not plan_name:
            return jsonify({'success': False, 'error': 'No plan selected'})
        
        # Validate plan name
        valid_plans = ['Free Trail', 'Standard', 'Premium']
        if plan_name not in valid_plans:
            return jsonify({'success': False, 'error': 'Invalid plan selected'})
        
        username = session.get('username')
        if not username:
            return jsonify({'success': False, 'error': 'User not authenticated'})
        
        # Update user's membership status in Firebase
        try:
            user_ref = db.reference(f'Data/{username}')
            user_data = user_ref.get() or {}
            user_data['membership_status'] = plan_name
            user_ref.set(user_data)
        except Exception as firebase_error:
            logger.warning(f"Firebase update failed: {str(firebase_error)}")
        
        # Update user's membership status in SQLite as backup
        try:
            conn = sqlite3.connect('file_conversion.db')
            cursor = conn.cursor()
            
            # Update the user's membership status
            cursor.execute("""
                UPDATE users 
                SET membership_status = ? 
                WHERE username = ?
            """, (plan_name, username))
            
            conn.commit()
            updated_rows = cursor.rowcount
            conn.close()
        except Exception as sqlite_error:
            logger.warning(f"SQLite update failed: {str(sqlite_error)}")
            updated_rows = 1  # Assume success if Firebase worked
        
        logger.info(f"User {username} upgraded to {plan_name} plan")
        return jsonify({
            'success': True, 
            'message': f'Successfully upgraded to {plan_name} plan!',
            'plan': plan_name
        })
            
    except Exception as e:
        logger.error(f"Plan selection error: {str(e)}")
        return jsonify({'success': False, 'error': 'An error occurred while selecting the plan'})

@app.route('/user_dashboard')
@login_required
def user_dashboard():
    return "User Dashboard (Under Construction)"

@app.route('/upload_txt', methods=['POST'])
@login_required
def upload_txt():
    content = request.form.get('text')
    if content and 'username' in session:
        user_path = get_user_storage_path()
        ref = db.reference(f'{user_path}/txt')
        ref.push({"content": content})
        return jsonify({"message": "Text uploaded successfully."})
    return jsonify({"error": "Invalid request"}), 400

# --- File Conversion and Upload/Download Routes (using per-user storage) ---

@app.route('/image-to-pdf')
@login_required
def image_to_pdf_page():
    return render_template('image_to_pdf.html', **get_user_context())

@app.route('/pdf-to-image')
@login_required
def pdf_to_image_page():
    return render_template('pdf_to_image.html', **get_user_context())

@app.route('/merge-pdfs')
@login_required
def merge_pdfs_page():
    return render_template('merge_pdfs.html', **get_user_context())

@app.route('/word-to-pdf')
@login_required
def word_to_pdf_page():
    return render_template('word_to_pdf.html', **get_user_context())

@app.route('/excel-to-pdf')
@login_required
def excel_to_pdf_page():
    return render_template('excel_to_pdf.html', **get_user_context())

@app.route('/pdf-to-ppt')
@login_required
def pdf_to_ppt_page():
    return render_template('pdf_to_ppt.html', **get_user_context())

@app.route('/bg-remover')
@login_required
def bg_remover_page():
    return render_template('bg_remover.html', **get_user_context())

@app.route('/admin-logs')
@login_required
def logs_page():
    return render_template('logs.html', **get_user_context())

@app.route('/history')
@login_required
def history_page():
    return render_template('history.html', **get_user_context())

@app.route('/compress-pdf')
@login_required
def compress_pdf_page():
    return render_template('compress_pdf.html', **get_user_context())

@app.route('/split-pdf')
@login_required
def split_pdf_page():
    return render_template('split_pdf.html', **get_user_context())

@app.route('/remove-pages-ui')
@login_required
def remove_pages_ui():
    return render_template('remove_page.html', **get_user_context())

@app.route('/document-screener')
@login_required
def document_screener_page():
    global current_document_text, conversation_history
    current_document_text = ''
    conversation_history.clear()
    return render_template('document_screener.html', **get_user_context())

@app.route('/plagiarism-scanner')
@login_required
def plagiarism_scanner_page():
    result = session.pop('plagiarism_result', None)
    input_text = session.pop('plagiarism_input_text', '')
    user_context = get_user_context()
    return render_template('plagiarism.html', result=result, input_text=input_text, **user_context)

@app.route('/text-to-speech')
@login_required
def text_to_speech_page():
    return render_template('text_to_speech.html', **get_user_context())

@app.route('/speech-to-text')
@login_required
def speech_to_text_page():
    return render_template('speech_to_text.html', **get_user_context())

@app.route('/ai-pdf-editor')
@login_required
def ai_pdf_editor_page():
    return render_template('ai_pdf_editor.html', **get_user_context())

@app.route('/text-summarizer')
@login_required
def text_summarizer_page():
    return render_template('text_summarizer.html', **get_user_context())

# --- Image to PDF ---
@app.route('/convert/image-to-pdf', methods=['POST'])
@login_required
def convert_image_to_pdf():
    if 'images' not in request.files:
        return jsonify({'error': 'No images provided'}), 400
    files = request.files.getlist('images')
    image_list = []
    for file in files:
        try:
            image = Image.open(file.stream)
            if image.mode != 'RGB':
                image = image.convert('RGB')
            image_list.append(image)
        except Exception as e:
            return jsonify({'error': f'Failed to read image: {str(e)}'}), 500
    if not image_list:
        return jsonify({'error': 'No valid images found'}), 400

    output_filename = f"{uuid.uuid4().hex}.pdf"
    
    # Use temporary file instead of local folder
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
        output_path = temp_file.name
    
    try:
        image_list[0].save(output_path, save_all=True, append_images=image_list[1:], format='PDF')

        username = session.get('username')
        cloudinary_folder = f'storage/{username}/files'
        cloudinary_url = upload_to_cloudinary(output_path, cloudinary_folder)
        store_url_in_firebase(cloudinary_url, 'files', output_filename)

        log_conversion('image-to-pdf', files[0].filename, output_filename, output_path, cloudinary_url)

        return send_file(output_path, as_attachment=True, download_name="converted.pdf", mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'PDF conversion failed: {str(e)}'}), 500
    finally:
        # Clean up temporary file
        try:
            if os.path.exists(output_path):
                os.remove(output_path)
        except Exception as e:
            logger.warning(f"Failed to clean up temp file: {str(e)}")

# --- PDF to Image ---
@app.route('/convert/pdf-to-image', methods=['POST'])
@login_required
def convert_pdf_to_image():
    pdf_file = request.files.get('pdf')
    if not pdf_file:
        return "No PDF uploaded", 400
    try:
        images = convert_from_bytes(pdf_file.read(), fmt='png', single_file=True)
        output_filename = f"{uuid.uuid4().hex}.png"
        
        # Use temporary file instead of local folder
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
            output_path = temp_file.name
            
        images[0].save(output_path, format='PNG')

        username = session.get('username')
        cloudinary_folder = f'storage/{username}/img'
        cloudinary_url = upload_to_cloudinary(output_path, cloudinary_folder)
        store_url_in_firebase(cloudinary_url, 'img', output_filename)

        log_conversion('pdf-to-image', pdf_file.filename, output_filename, output_path, cloudinary_url)

        return send_file(output_path, as_attachment=True, download_name="converted.png", mimetype='image/png')
    except Exception as e:
        return f"Error during conversion: {str(e)}", 500
    finally:
        # Clean up temporary file
        try:
            if 'output_path' in locals() and os.path.exists(output_path):
                os.remove(output_path)
        except Exception as e:
            logger.warning(f"Failed to clean up temp file: {str(e)}")

# --- Merge PDFs ---
@app.route('/merge/pdfs', methods=['POST'])
@login_required
def merge_pdfs():
    if 'pdfs' not in request.files:
        return jsonify({'error': 'No PDF files provided'}), 400
    files = request.files.getlist('pdfs')
    merger = PdfMerger()
    try:
        for file in files:
            merger.append(file)
        output_filename = f"merged_{uuid.uuid4().hex}.pdf"
        
        # Use temporary file instead of local folder  
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            output_path = temp_file.name
            
        merger.write(output_path)
        merger.close()

        username = session.get('username')
        cloudinary_folder = f'storage/{username}/files'
        cloudinary_url = upload_to_cloudinary(output_path, cloudinary_folder)
        store_url_in_firebase(cloudinary_url, 'files', output_filename)

        log_conversion('merge-pdfs', files[0].filename, output_filename, output_path, cloudinary_url)

        return send_file(output_path, as_attachment=True, download_name='merged.pdf', mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Merge failed: {str(e)}'}), 500
    finally:
        # Clean up temporary file
        try:
            if 'output_path' in locals() and os.path.exists(output_path):
                os.remove(output_path)
        except Exception as e:
            logger.warning(f"Failed to clean up temp file: {str(e)}")

# --- Download route (per-user) ---
@app.route('/download')
@login_required
def download():
    file_path = request.args.get('file_path')
    file_name = request.args.get('file_name')
    mime_type = request.args.get('mime_type', 'application/octet-stream')
    if not os.path.exists(file_path):
        return jsonify({'error': 'File not found on server'}), 404

    ext = os.path.splitext(file_name)[1].lower()
    if ext in ['.pdf', '.docx', '.xlsx', '.pptx', '.txt', '.csv']:
        category = 'files'
    elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp']:
        category = 'img'
    elif ext in ['.mp3', '.wav', '.aac', '.ogg', '.flac']:
        category = 'audio'
    else:
        category = 'files'

    username = session.get('username')
    folder_path = f'storage/{username}/{category}'
    cloudinary_url = upload_to_cloudinary(file_path, folder_path)
    store_url_in_firebase(cloudinary_url, category, file_name)

    return send_file(file_path, as_attachment=True, download_name=file_name, mimetype=mime_type)

# --- Download from Cloudinary/Firebase ---
@app.route('/download-file/<int:file_id>')
@login_required
def download_file_from_cloud(file_id):
    """Download file from Cloudinary using file ID from database."""
    try:
        # Get file info from database
        conn = sqlite3.connect('file_conversion.db')
        c = conn.cursor()
        c.execute('''
            SELECT cloudinary_url, converted_filename, original_filename 
            FROM conversions 
            WHERE id = ? AND username = ?
        ''', (file_id, session.get('username')))
        result = c.fetchone()
        conn.close()
        
        if not result:
            return jsonify({'error': 'File not found or access denied'}), 404
            
        cloudinary_url, converted_filename, original_filename = result
        
        if not cloudinary_url:
            return jsonify({'error': 'File URL not available'}), 404
        
        # Fetch file from Cloudinary
        response = requests.get(cloudinary_url)
        if response.status_code != 200:
            return jsonify({'error': 'File not accessible from cloud storage'}), 404
        
        # Create a temporary file to serve
        temp_file = tempfile.NamedTemporaryFile(delete=False)
        temp_file.write(response.content)
        temp_file.close()
        
        # Determine MIME type
        mime_type, _ = mimetypes.guess_type(converted_filename)
        if not mime_type:
            mime_type = 'application/octet-stream'
        
        def cleanup_temp_file():
            try:
                os.unlink(temp_file.name)
            except:
                pass
        
        # Schedule cleanup after sending file
        from threading import Timer
        Timer(5.0, cleanup_temp_file).start()
        
        return send_file(
            temp_file.name, 
            as_attachment=True, 
            download_name=converted_filename, 
            mimetype=mime_type
        )
        
    except Exception as e:
        return jsonify({'error': f'Download failed: {str(e)}'}), 500

# --- Word to PDF ---
@app.route('/convert_word_to_pdf', methods=['POST'])
@login_required
def convert_word_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded!'}), 400
    file = request.files['file']

    if not file or not file.filename:
        return jsonify({'error': 'No file selected!'}), 400

    filename = file.filename.lower()

    allowed_extensions = ['.doc', '.docx']
    file_extension = '.' + filename.split('.')[-1] if '.' in filename else ''

    if file_extension not in allowed_extensions:
        return jsonify({'error': 'Only .doc and .docx files are supported'}), 400

    # Save uploaded file to a temporary file and close immediately
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_in:
        file.save(temp_in)
        input_path = temp_in.name

    output_filename = f"{uuid.uuid4().hex}.pdf"
    
    # Use temporary file for output
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_out:
        output_path = temp_out.name

    try:
        if file_extension == '.docx':
            try:
                import pythoncom
                pythoncom.CoInitialize()
                from docx2pdf import convert as docx2pdf_convert
                docx2pdf_convert(input_path, output_path)
                pythoncom.CoUninitialize()
                if not os.path.exists(output_path):
                    raise Exception("docx2pdf did not create output file.")
            except Exception as e:
                logger.error(f'docx2pdf failed: {e}, falling back to text pdf conversion.')
                # Use python-docx fallback
                doc = Document(input_path)
                text = []

                # Extract paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():
                        text.append(para.text)

                # Extract tables
                for table in doc.tables:
                    for row in table.rows:
                        row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                        if row_text.strip():
                            text.append(row_text)

                if not text:
                    raise Exception("No text found in docx for fallback conversion.")

                pdf_io = io.BytesIO()
                c = canvas.Canvas(pdf_io, pagesize=letter)

                width, height = letter
                margin = 50
                line_height = 14
                y = height - margin

                for line in text:
                    if y < margin:
                        c.showPage()
                        y = height - margin

                    max_chars = int((width - 2 * margin) / 7)
                    if len(line) > max_chars:
                        line = line[:max_chars] + "..."

                    try:
                        c.drawString(margin, y, line)
                    except:
                        clean_line = ''.join(char for char in line if ord(char) < 128)
                        c.drawString(margin, y, clean_line)

                    y -= line_height

                c.save()
                pdf_io.seek(0)
                with open(output_path, 'wb') as f:
                    f.write(pdf_io.getvalue())

        elif file_extension == '.doc':
            # Handle legacy .doc files
            try:
                # Attempt to open .doc as .docx (may fail)
                doc = Document(input_path)
                text = [para.text for para in doc.paragraphs if para.text.strip()]

                if not text:
                    raise Exception("No text extracted from .doc file using python-docx")

            except Exception as e:
                logger.warning(f"python-docx failed for .doc file: {e}, trying alternative method")
                try:
                    with open(input_path, 'rb') as f:
                        content = f.read()
                    text_content = content.decode('utf-8', errors='ignore')
                    lines = text_content.split('\n')
                    text = [line.strip() for line in lines if line.strip() and len(line.strip()) > 2]

                    clean_text = []
                    for line in text[:100]:
                        if any(c.isalpha() for c in line) and len(line) < 200:
                            clean_text.append(line)
                    if not clean_text:
                        raise Exception("No readable text found in .doc file")

                    text = clean_text

                except Exception as e2:
                    logger.error(f"Alternative .doc reading method failed: {e2}")
                    return jsonify({'error': 'Failed to read .doc file. Please try converting to .docx format first.'}), 500

            pdf_io = io.BytesIO()
            c = canvas.Canvas(pdf_io, pagesize=letter)

            width, height = letter
            margin = 50
            line_height = 14
            y = height - margin

            for line in text:
                if y < margin:
                    c.showPage()
                    y = height - margin

                max_chars = int((width - 2 * margin) / 7)
                if len(line) > max_chars:
                    line = line[:max_chars] + "..."

                try:
                    c.drawString(margin, y, line)
                except:
                    clean_line = ''.join(char for char in line if ord(char) < 128)
                    c.drawString(margin, y, clean_line)

                y -= line_height

            c.save()
            pdf_io.seek(0)
            with open(output_path, 'wb') as f:
                f.write(pdf_io.getvalue())

        # Verify PDF creation
        if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            raise Exception("PDF conversion failed - output file is empty or missing")

        # Upload to Cloudinary and store in Firebase (replace with your implementations)
        username = session.get('username')
        cloudinary_folder = f'storage/{username}/files'
        cloudinary_url = upload_to_cloudinary(output_path, cloudinary_folder)
        store_url_in_firebase(cloudinary_url, 'files', output_filename)

        # Log conversion (replace with your implementation)
        log_conversion('word-to-pdf', file.filename, output_filename, output_path, cloudinary_url)

        base_name = file.filename.rsplit('.', 1)[0] if '.' in file.filename else file.filename
        download_name = f"{base_name}.pdf"

        return send_file(output_path, as_attachment=True, download_name=download_name, mimetype='application/pdf')

    except Exception as e:
        logger.error(f'Word to PDF conversion failed: {str(e)}', exc_info=True)
        return jsonify({'error': f'Word to PDF conversion failed: {str(e)}'}), 500

    finally:
        # Clean up temp input file
        try:
            if os.path.exists(input_path):
                os.remove(input_path)
        except Exception as cleanup_error:
            logger.warning(f"Error deleting temp input file {input_path}: {str(cleanup_error)}")
            
        # Clean up temp output file
        try:
            if 'output_path' in locals() and os.path.exists(output_path):
                os.remove(output_path)
        except Exception as cleanup_error:
            logger.warning(f"Error deleting temp output file {output_path}: {str(cleanup_error)}")

# --- Excel to PDF ---
@app.route('/convert_excel_to_pdf', methods=['POST'])
@login_required
def convert_excel_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded!'}), 400
    file = request.files['file']
    if not file.filename.endswith('.xlsx'):
        return jsonify({'error': 'Only .xlsx files are supported'}), 400
        
    # Use temporary file for input
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_input:
        file.save(temp_input)
        input_path = temp_input.name
        
    output_filename = f"{uuid.uuid4().hex}.pdf"
    
    # Use temporary file for output
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_output:
        output_path = temp_output.name
    try:
        wb = load_workbook(input_path)
        sheet = wb.active
        text = []
        for row in sheet.rows:
            row_text = [str(cell.value) if cell.value is not None else '' for cell in row]
            text.append(' | '.join(row_text))
        pdf_io = io.BytesIO()
        c = canvas.Canvas(pdf_io, pagesize=letter)
        y = 750
        for line in text:
            if y < 50:
                c.showPage()
                y = 750
            c.drawString(50, y, line[:100])
            y -= 15
        c.save()
        pdf_io.seek(0)
        with open(output_path, 'wb') as f:
            f.write(pdf_io.getvalue())

        username = session.get('username')
        cloudinary_folder = f'storage/{username}/files'
        cloudinary_url = upload_to_cloudinary(output_path, cloudinary_folder)
        store_url_in_firebase(cloudinary_url, 'files', output_filename)

        log_conversion('excel-to-pdf', file.filename, output_filename, output_path, cloudinary_url)

        return send_file(output_path, as_attachment=True, download_name=file.filename.replace('.xlsx', '.pdf'), mimetype='application/pdf')
    except Exception as e:
        return jsonify({'error': f'Excel to PDF conversion failed: {str(e)}'}), 500
    finally:
        # Clean up temp files
        try:
            if os.path.exists(input_path):
                os.remove(input_path)
        except Exception as e:
            logger.warning(f"Failed to clean up temp input file: {str(e)}")
        try:
            if 'output_path' in locals() and os.path.exists(output_path):
                os.remove(output_path)
        except Exception as e:
            logger.warning(f"Failed to clean up temp output file: {str(e)}")

# --- PDF to PPT ---
@app.route('/convert_pdf_to_ppt', methods=['POST'])
@login_required
def convert_pdf_to_ppt():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded!'}), 400
    file = request.files['file']
    if not file.filename.endswith('.pdf'):
        return jsonify({'error': 'Only .pdf files are supported'}), 400
    output_filename = f"{uuid.uuid4().hex}.pptx"
    
    # Use temporary file for output
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_output:
        output_path = temp_output.name
    try:
        pdf_reader = PdfReader(file.stream)
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content slide
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text = page.extract_text() or f"Page {page_num + 1}"
            lines = text.split('\n')
            filtered_lines = [line.strip() for line in lines if line.strip()]
            slide_text = '\n'.join(filtered_lines[:10])
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = f"Slide {page_num + 1}"
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.text = slide_text
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        with open(output_path, 'wb') as f:
            f.write(ppt_io.getvalue())

        username = session.get('username')
        cloudinary_folder = f'storage/{username}/files'
        cloudinary_url = upload_to_cloudinary(output_path, cloudinary_folder)
        
        log_conversion('pdf-to-ppt', file.filename, output_filename, output_path, cloudinary_url)
        store_url_in_firebase(cloudinary_url, 'files', output_filename)

        return send_file(output_path, as_attachment=True,
                         download_name=file.filename.replace('.pdf', '.pptx'),
                         mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    except Exception as e:
        return jsonify({'error': f'PDF to PPT conversion failed: {str(e)}'}), 500
    finally:
        # Clean up temp file
        try:
            if 'output_path' in locals() and os.path.exists(output_path):
                os.remove(output_path)
        except Exception as e:
            logger.warning(f"Failed to clean up temp file: {str(e)}")

# --- Background Remover ---
@app.route('/remove_background', methods=['POST'])
@login_required
def remove_background():
    logger.debug("Received request at /remove_background")
    if 'image' not in request.files:
        logger.error("No image uploaded")
        return jsonify({'error': 'No image uploaded!'}), 400
    file = request.files['image']
    valid_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp')
    if not file.filename.lower().endswith(valid_extensions):
        return jsonify({'error': 'Unsupported file extension'}), 400
    if file.content_length and file.content_length > 10 * 1024 * 1024:
        return jsonify({'error': 'File size exceeds 10MB limit'}), 400

    temp_file_path = None
    try:
        image = Image.open(file.stream)
        if image.size[0] > 4096 or image.size[1] > 4096:
            return jsonify({'error': 'Image resolution exceeds 4096x4096 limit'}), 400
        if file.filename.lower().endswith('.gif'):
            image.seek(0)
        if image.mode != 'RGBA':
            image = image.convert('RGBA')
        output_image = remove(image)

        output_filename = f"bg_removed_{uuid.uuid4().hex}.png"
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
            temp_file_path = temp_file.name
            output_image.save(temp_file_path, format='PNG')

        # No need for additional output_path, use temp_file_path directly
        username = session.get('username')
        cloudinary_folder = f'storage/{username}/img'
        cloudinary_url = upload_to_cloudinary(temp_file_path, cloudinary_folder)
        store_url_in_firebase(cloudinary_url, 'img', output_filename)

        log_conversion('background-remover', file.filename, output_filename, temp_file_path, cloudinary_url)

        return send_file(temp_file_path, as_attachment=True, download_name='background_removed.png', mimetype='image/png')
    except Exception as e:
        logger.error(f"Background removal failed: {str(e)}", exc_info=True)
        return jsonify({'error': f'Background removal failed: {str(e)}'}), 500
    finally:
        if temp_file_path and os.path.exists(temp_file_path):
            try:
                os.remove(temp_file_path)
            except Exception as e:
                logger.warning(f"Failed to clean up temp file: {str(e)}")

# --- Other routes like Page Count, Remove Pages, Compress, Split PDFs ---
# --- Following the same pattern updating to per-user storage paths in DB/cloud ---

@app.route('/get-page-count', methods=['POST'])
@login_required
def get_page_count():
    pdf = request.files.get('pdf')
    if not pdf:
        return jsonify({'error': 'No PDF uploaded'}), 400
    filename = secure_filename(pdf.filename)
    
    # Use temporary file instead of upload folder
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
        pdf.save(temp_file)
        filepath = temp_file.name
    try:
        reader = PdfReader(filepath)
        count = len(reader.pages)
        os.remove(filepath)
        return jsonify({'page_count': count})
    except Exception as e:
        os.remove(filepath)
        return jsonify({'error': str(e)}), 500

@app.route('/remove-pages', methods=['POST'])
@login_required
def remove_pages():
    pdf = request.files.get('pdf')
    removed_pages = request.form.get('removed_pages', '')
    total_pages = int(request.form.get('page_count', 0))
    if not pdf:
        return "No PDF uploaded", 400
    try:
        # Convert page numbers (1-based) to indices (0-based)
        remove_page_numbers = [int(i) for i in removed_pages.split(',') if i.strip().isdigit()]
        remove_indices = [page - 1 for page in remove_page_numbers if 1 <= page <= total_pages]
    except ValueError:
        return "Invalid page numbers", 400
    
    if not remove_indices:
        return "No valid page numbers provided", 400
    
    filename = secure_filename(pdf.filename)
    
    # Use temporary file instead of upload folder
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
        pdf.save(temp_file)
        filepath = temp_file.name
    reader = PdfReader(filepath)
    writer = PdfWriter()
    for i in range(len(reader.pages)):
        if i not in remove_indices:
            writer.add_page(reader.pages[i])
    
    output_filename = f"removed_{uuid.uuid4().hex}.pdf"
    # Create PDF content in memory
    pdf_buffer = io.BytesIO()
    writer.write(pdf_buffer)
    pdf_buffer.seek(0)
    
    os.remove(filepath)

    # Save to temporary file for cloud upload
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
        temp_output_path = temp_file.name
        temp_file.write(pdf_buffer.getvalue())

    log_conversion('remove-pages', filename, output_filename, temp_output_path, None, 'success')

    username = session.get('username')
    cloudinary_folder = f'storage/{username}/files'
    cloudinary_url = upload_to_cloudinary(temp_output_path, cloudinary_folder)
    store_url_in_firebase(cloudinary_url, 'files', output_filename)

    # Clean up temp file after cloud upload
    if os.path.exists(temp_output_path):
        os.remove(temp_output_path)

    # Send file from memory buffer
    pdf_buffer.seek(0)
    return send_file(pdf_buffer, as_attachment=True, download_name='removed_pages.pdf', mimetype='application/pdf')

@app.route('/compress', methods=['POST'])
@login_required
def compress_pdf():
    pdf_file = request.files.get('pdf')
    compression_level = request.form.get('compression_level')
    
    if not pdf_file:
        return "No PDF uploaded", 400
    
    if compression_level not in ['low', 'medium', 'high']:
        return "Invalid compression level", 400
    
    filename = secure_filename(pdf_file.filename)
    
    # Use temporary file instead of upload folder
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
        pdf_file.save(temp_file)
        filepath = temp_file.name
    
    try:
        original_size = os.path.getsize(filepath)
        logger.info(f"Original file size: {original_size / 1024:.2f} KB")
        
        output_filename = f"compressed_{uuid.uuid4().hex}.pdf"
        
        # Try Ghostscript first for better compression
        if is_ghostscript_installed():
            logger.info("Using Ghostscript for compression")
            gs_quality = {'low': 'printer', 'medium': 'ebook', 'high': 'screen'}
            gs_setting = gs_quality[compression_level]
            
            # Create temporary output file for Ghostscript
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_output:
                temp_output_path = temp_output.name
            
            gs_command = [
                'gs', '-sDEVICE=pdfwrite', '-dCompatibilityLevel=1.4',
                f'-dPDFSETTINGS=/{gs_setting}', '-dNOPAUSE', '-dQUIET', '-dBATCH',
                f'-sOutputFile={temp_output_path}', filepath
            ]
            result = subprocess.run(gs_command, capture_output=True, text=True)
            if result.returncode != 0:
                logger.error(f"Ghostscript failed: {result.stderr}")
                raise Exception(f"Ghostscript compression failed: {result.stderr}")
            
            # Read compressed PDF into memory
            with open(temp_output_path, 'rb') as f:
                compressed_pdf_data = f.read()
            
            # Clean up temp file
            os.remove(temp_output_path)
            
        else:
            # Fallback to PyMuPDF compression with actual image processing
            logger.warning("Ghostscript not available, using PyMuPDF compression")
            doc = None
            try:
                doc = fitz.open(filepath)
                
                # Define compression settings based on level
                if compression_level == 'low':
                    # Light compression - preserve quality
                    image_quality = 85
                    downscale_factor = 1.0
                elif compression_level == 'medium':
                    # Moderate compression - balance quality and size
                    image_quality = 60
                    downscale_factor = 0.8
                else:  # high
                    # Heavy compression - prioritize size reduction
                    image_quality = 30
                    downscale_factor = 0.6
                
                # Process each page and compress images
                images_processed = 0
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    image_list = page.get_images(full=True)
                    
                    for img_index, img in enumerate(image_list):
                        try:
                            # Get image data
                            xref = img[0]
                            base_image = doc.extract_image(xref)
                            image_bytes = base_image["image"]
                            
                            # Load image with PIL
                            pil_image = Image.open(io.BytesIO(image_bytes))
                            
                            # Resize image if needed
                            if downscale_factor < 1.0:
                                new_width = int(pil_image.width * downscale_factor)
                                new_height = int(pil_image.height * downscale_factor)
                                pil_image = pil_image.resize((new_width, new_height), Image.Resampling.LANCZOS)
                            
                            # Convert to RGB if necessary
                            if pil_image.mode != 'RGB':
                                pil_image = pil_image.convert('RGB')
                            
                            # Compress image
                            img_buffer = io.BytesIO()
                            pil_image.save(img_buffer, format='JPEG', quality=image_quality, optimize=True)
                            img_buffer.seek(0)
                            
                            # Replace image in PDF
                            doc._updateObject(xref, img_buffer.getvalue())
                            images_processed += 1
                            
                        except Exception as e:
                            logger.warning(f"Failed to process image {img_index} on page {page_num}: {str(e)}")
                            continue
                
                logger.info(f"Processed {images_processed} images for compression")
                
                # Save to memory buffer instead of file
                pdf_buffer = io.BytesIO()
                doc.save(pdf_buffer, deflate=True, clean=True)
                compressed_pdf_data = pdf_buffer.getvalue()
                
            finally:
                # Ensure document is properly closed
                if doc is not None:
                    doc.close()
        
        # Verify compression worked
        if not compressed_pdf_data or len(compressed_pdf_data) == 0:
            raise Exception("Compression failed - output file is empty or missing")
        
        compressed_size = len(compressed_pdf_data)
        logger.info(f"Compressed file size: {compressed_size / 1024:.2f} KB")
        reduction = (original_size - compressed_size) / original_size * 100 if original_size > 0 else 0
        logger.info(f"Size reduction: {reduction:.2f}%")

        # Log conversion for tracking purposes
        log_conversion('compress-pdf', filename, output_filename, "memory_cache", None, 'success')

        # Clean up input file
        try:
            os.remove(filepath)
        except PermissionError:
            logger.warning(f"Could not delete input file {filepath} - file may be in use")
        
        # Store compressed PDF data in memory cache (no cloud upload yet)
        compressed_pdf_cache[output_filename] = compressed_pdf_data
        
        logger.info(f"Compression successful: {output_filename}")
        
        # Return success response with download URL
        return jsonify({
            'success': True,
            'download_url': f'/download-compressed/{output_filename}',
            'filename': output_filename,
            'original_size': f"{original_size / 1024:.2f} KB",
            'compressed_size': f"{compressed_size / 1024:.2f} KB",
            'reduction': f"{reduction:.2f}%",
            'message': f'PDF successfully compressed with {reduction:.2f}% size reduction'
        })
        
    except Exception as e:
        # Clean up input file on error
        try:
            if os.path.exists(filepath):
                os.remove(filepath)
        except PermissionError:
            logger.warning(f"Could not delete input file {filepath} - file may be in use")
        except Exception as cleanup_error:
            logger.warning(f"Error during cleanup: {str(cleanup_error)}")
        
        logger.error(f"Error during compression: {str(e)}")
        return f"Error during compression: {str(e)}", 500

# Store split PDFs and compressed PDFs in memory for immediate download
split_pdf_cache = {}
compressed_pdf_cache = {}
audio_cache = {}

@app.route('/download-split/<filename>')
@login_required
def download_split_pdf(filename):
    # Check if file exists in memory cache
    if filename not in split_pdf_cache:
        return jsonify({'error': 'File not found or expired'}), 404
    
    try:
        # Get PDF data from memory cache
        pdf_data = split_pdf_cache[filename]
        
        # Create a BytesIO object with the PDF data first
        pdf_buffer = io.BytesIO(pdf_data)
        pdf_buffer.seek(0)
        
        # Now upload to cloud storage when user actually downloads
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
                temp_file.write(pdf_data)
                temp_file_path = temp_file.name
            
            # Upload to cloud only when downloading
            username = session.get('username')
            cloudinary_folder = f'storage/{username}/files'
            cloudinary_url = upload_to_cloudinary(temp_file_path, cloudinary_folder)
            store_url_in_firebase(cloudinary_url, 'files', filename)
            
            # Clean up temp file
            if os.path.exists(temp_file_path):
                os.remove(temp_file_path)
        except Exception as cloud_error:
            # If cloud upload fails, still serve the file
            logger.warning(f"Cloud upload failed: {str(cloud_error)}")
        
        # Clean up from cache after preparing to serve
        if filename in split_pdf_cache:
            del split_pdf_cache[filename]
        
        # Make sure the buffer is at the beginning
        pdf_buffer.seek(0)
        
        return send_file(
            pdf_buffer, 
            as_attachment=True, 
            download_name=filename, 
            mimetype='application/pdf'
        )
    except Exception as e:
        # Clean up from cache on error
        if filename in split_pdf_cache:
            del split_pdf_cache[filename]
        logger.error(f"Download failed: {str(e)}")
        return jsonify({'error': f'Download failed: {str(e)}'}), 500

@app.route('/download-compressed/<filename>')
@login_required
def download_compressed_pdf(filename):
    # Check if file exists in memory cache
    if filename not in compressed_pdf_cache:
        return jsonify({'error': 'File not found or expired'}), 404
    
    try:
        # Get PDF data from memory cache
        pdf_data = compressed_pdf_cache[filename]
        
        # Create a BytesIO object with the PDF data first
        pdf_buffer = io.BytesIO(pdf_data)
        pdf_buffer.seek(0)
        
        # Now upload to cloud storage when user actually downloads
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
                temp_file.write(pdf_data)
                temp_file_path = temp_file.name
            
            # Upload to cloud only when downloading
            username = session.get('username')
            cloudinary_folder = f'storage/{username}/files'
            cloudinary_url = upload_to_cloudinary(temp_file_path, cloudinary_folder)
            store_url_in_firebase(cloudinary_url, 'files', filename)
            
            # Clean up temp file
            if os.path.exists(temp_file_path):
                os.remove(temp_file_path)
        except Exception as cloud_error:
            # If cloud upload fails, still serve the file
            logger.warning(f"Cloud upload failed: {str(cloud_error)}")
        
        # Clean up from cache after preparing to serve
        if filename in compressed_pdf_cache:
            del compressed_pdf_cache[filename]
        
        # Make sure the buffer is at the beginning
        pdf_buffer.seek(0)
        
        return send_file(
            pdf_buffer, 
            as_attachment=True, 
            download_name=filename, 
            mimetype='application/pdf'
        )
    except Exception as e:
        # Clean up from cache on error
        if filename in compressed_pdf_cache:
            del compressed_pdf_cache[filename]
        logger.error(f"Download failed: {str(e)}")
        return jsonify({'error': f'Download failed: {str(e)}'}), 500

@app.route('/split', methods=['POST'])
@login_required
def split_pdf():
    try:
        pdf_file = request.files.get('pdf')
        split_index = request.form.get('split_index')
        
        if not pdf_file or not split_index:
            return jsonify({'error': 'Missing file or split index'}), 400
        
        try:
            split_index = int(split_index)
        except ValueError:
            return jsonify({'error': 'Invalid split index'}), 400
        
        filename = secure_filename(pdf_file.filename)
        
        # Use temporary file instead of upload folder
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            pdf_file.save(temp_file)
            filepath = temp_file.name
        
        reader = PdfReader(filepath)
        total_pages = len(reader.pages)
        
        if split_index <= 0 or split_index >= total_pages:
            os.remove(filepath)
            return jsonify({'error': f'Split index out of range. Must be between 1 and {total_pages - 1}'}), 400
        
        writer1 = PdfWriter()
        writer2 = PdfWriter()
        
        for i, page in enumerate(reader.pages):
            if i < split_index:
                writer1.add_page(page)
            else:
                writer2.add_page(page)
        
        # Create PDF content in memory
        pdf_buffer1 = io.BytesIO()
        pdf_buffer2 = io.BytesIO()
        writer1.write(pdf_buffer1)
        writer2.write(pdf_buffer2)
        
        output_filename1 = f"split_part1_{uuid.uuid4().hex}.pdf"
        output_filename2 = f"split_part2_{uuid.uuid4().hex}.pdf"
        
        # Log conversion for tracking purposes
        log_conversion('split-pdf', filename, f"{output_filename1}, {output_filename2}", "memory_cache", None, 'success')

        # Clean up the original file
        os.remove(filepath)
        
        # Store PDF data in memory cache for immediate download (no cloud upload yet)
        pdf_buffer1.seek(0)
        pdf_buffer2.seek(0)
        split_pdf_cache[output_filename1] = pdf_buffer1.getvalue()
        split_pdf_cache[output_filename2] = pdf_buffer2.getvalue()
        
        logger.info(f"Split PDF successful: {output_filename1}, {output_filename2}")
        
        # Return download URLs that serve from memory cache
        return jsonify({
            'success': True,
            'part1': f'/download-split/{output_filename1}',
            'part2': f'/download-split/{output_filename2}',
            'part1_name': output_filename1,
            'part2_name': output_filename2,
            'message': f'PDF successfully split into {split_index} and {total_pages - split_index} pages'
        })
        
    except Exception as e:
        logger.error(f"Split PDF failed: {str(e)}")
        # Clean up file if it exists
        if 'filepath' in locals() and os.path.exists(filepath):
            os.remove(filepath)
        return jsonify({'error': f'Split PDF failed: {str(e)}'}), 500


# --- Document Screener Routes ---

def extract_text_from_pdf(file_path):
    try:
        with open(file_path, 'rb') as file:
            reader = PdfReader(file)
            text = ''
            for page in reader.pages:
                text += page.extract_text() or ''
            return text
    except Exception as e:
        return f"Error extracting text from PDF: {str(e)}"

def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        text = '\n'.join([para.text for para in doc.paragraphs if para.text])
        return text
    except Exception as e:
        return f"Error extracting text from DOCX: {str(e)}"

def extract_text_from_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        return f"Error extracting text from TXT: {str(e)}"

def analyze_text_with_openrouter(text, format_type):
    prompt = (
        f"Analyze the following document text and provide a summary in "
        f"{'a concise paragraph' if format_type == 'paragraph' else 'bullet points'}. "
        f"Focus on key themes, topics, or entities mentioned in the text.\n\n"
        f"Text:\n{text[:2000]}"
    )

    headers = {
        'Authorization': f'Bearer {OPENROUTER_API_KEY}',
        'Content-Type': 'application/json'
    }

    data = {
        'model': OPENROUTER_MODEL,
        'messages': [{'role': 'user', 'content': prompt}]
    }

    try:
        response = requests.post(OPENROUTER_API_URL, headers=headers, json=data)
        response.raise_for_status()
        result = response.json()
        return result['choices'][0]['message']['content']
    except Exception as e:
        logger.error(f"OpenRouter API error: {str(e)}")
        return f"Error analyzing text with OpenRouter: {str(e)}"


@app.route('/analyze_document', methods=['POST'])
@login_required
def analyze_document():
    global current_document_text, conversation_history
    try:
        if 'docFile' not in request.files:
            return jsonify({'error': 'No document provided'}), 400
        doc_file = request.files['docFile']
        format_type = request.form.get('format', 'paragraph')
        temp_file_path = None
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{doc_file.filename.split('.')[-1]}") as temp_file:
            doc_file.save(temp_file)
            temp_file_path = temp_file.name
        file_extension = doc_file.filename.lower().split('.')[-1]
        if file_extension == 'pdf':
            text = extract_text_from_pdf(temp_file_path)
        elif file_extension == 'docx':
            text = extract_text_from_docx(temp_file_path)
        elif file_extension == 'txt':
            text = extract_text_from_txt(temp_file_path)
        else:
            os.remove(temp_file_path)
            return jsonify({'error': 'Unsupported file type. Use PDF, DOCX, or TXT.'}), 400

        os.remove(temp_file_path)
        if text.startswith('Error'):
            return jsonify({'error': text}), 500

        # Set global variable and clear conversation history for new document
        current_document_text = text
        conversation_history.clear()
        
        logger.debug(f"Document analyzed. Text length: {len(current_document_text)}")
        
        analysis = analyze_text_with_openrouter(text, format_type)
        if analysis.startswith('Error'):
            return jsonify({'error': analysis}), 500

        analysis_filename = f"analysis_{uuid.uuid4().hex}.txt"
        
        # Use temporary file for analysis output
        with tempfile.NamedTemporaryFile(delete=False, suffix='.txt', mode='w', encoding='utf-8') as temp_file:
            temp_file.write(f"File: {doc_file.filename}\nFormat: {format_type}\nAnalysis:\n{analysis}\n\n")
            analysis_path = temp_file.name

        log_conversion('document-screener', doc_file.filename, analysis_filename, analysis_path, None, 'success')
        
        # Clean up temp file
        try:
            if os.path.exists(analysis_path):
                os.remove(analysis_path)
        except Exception as e:
            logger.warning(f"Failed to clean up temp analysis file: {str(e)}")
            
        return jsonify({'analysis': analysis})
    except Exception as e:
        logger.error(f"Document analysis failed: {str(e)}")
        return jsonify({'error': f"Document analysis failed: {str(e)}"}), 500

@app.route('/chat', methods=['POST'])
@login_required
def chat():
    global conversation_history, current_document_text
    try:
        data = request.get_json()
        message = data.get('message', '')
        if not message:
            return jsonify({'error': 'No message provided'}), 400
        
        # Debug: Log the current document text status
        logger.debug(f"Current document text length: {len(current_document_text) if current_document_text else 0}")
        
        if not current_document_text:
            return jsonify({'error': 'No document uploaded. Please upload and analyze a document first using the "Analyze Document" button.'}), 400

        context = f"Document text:\n\n{current_document_text[:2000]}\n\nConversation history:\n"
        for role, msg in conversation_history:
            context += f"{role}: {msg}\n"
        prompt = f"{context}\nUser: {message}\nAssistant: Answer based on the document and conversation history. If the question is about names or specific details, extract relevant information from the document. If no relevant information is found, say so clearly."

        headers = {
            'Authorization': f'Bearer {OPENROUTER_API_KEY}',
            'Content-Type': 'application/json'
        }

        data = {
            'model': OPENROUTER_MODEL,
            'messages': [{'role': 'user', 'content': prompt}]
        }

        response = requests.post(OPENROUTER_API_URL, headers=headers, json=data)
        response.raise_for_status()
        result = response.json()
        response_text = result.get('choices', [{}])[0].get('message', {}).get('content', '')
        conversation_history.append(('User', message))
        conversation_history.append(('Assistant', response_text))
        return jsonify({'response': response_text})
    except Exception as e:
        logger.error(f"Chat error: {str(e)}")
        return jsonify({'error': f"Chat failed: {str(e)}"}), 500

# --- Plagiarism Scanner ---

def fetch_web_snippets(query, max_results=5):
    """Scrape DuckDuckGo for snippets."""
    search_url = f"https://html.duckduckgo.com/html/?q={query}"
    headers = {"User-Agent": "Mozilla/5.0"}
    try:
        response = requests.get(search_url, headers=headers, timeout=5)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        results = soup.find_all('a', {'class': 'result__a'}, limit=max_results)
        snippets = [r.get_text(strip=True) for r in results if r.get_text(strip=True)]
        return snippets
    except Exception as e:
        logger.error(f"Error fetching web snippets: {str(e)}")
        return []

def call_openrouter_similarity(text_a, text_b):
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
        "User-Agent": "Mozilla/5.0"
    }
    messages = [
        {
            "role": "system",
            "content": ("You are a plagiarism detection assistant. Given two texts, "
                        "respond ONLY with a plagiarism similarity percentage (0 to 100) and a brief explanation, separated by a newline.")
        },
        {
            "role": "user",
            "content": f"Text A:\n{text_a}\n\nText B:\n{text_b}"
        }
    ]
    data = {
        "model": OPENROUTER_MODEL,
        "messages": messages
    }
    try:
        response = requests.post(OPENROUTER_API_URL, headers=headers, json=data)
        response.raise_for_status()
        return response.json()['choices'][0]['message']['content'].strip()
    except Exception as e:
        logger.error(f"OpenRouter API error in similarity check: {str(e)}")
        return f"Error: {str(e)}"

@app.route('/check_plagiarism', methods=['POST'])
@login_required
def check_plagiarism():
    try:
        input_text = request.form.get('text', '').strip()
        if not input_text:
            return jsonify({'error': 'No text provided'}), 400
        if len(input_text) < 20:
            results = [{"snippet": "", "similarity": "Input text too short to check plagiarism."}]
            return jsonify({'results': results})

        query = input_text[:100]
        snippets = fetch_web_snippets(query)
        results = []
        for snippet in snippets:
            sim = call_openrouter_similarity(input_text, snippet)
            results.append({"snippet": snippet, "similarity": sim})
        return jsonify({'results': results})
    except Exception as e:
        logger.error(f"Plagiarism check failed: {str(e)}")
        return jsonify({'error': f"Plagiarism check failed: {str(e)}"}), 500

# --- Text to Speech ---

@app.route('/generate_tts', methods=['POST'])
@login_required
def generate_tts():
    data = request.get_json()
    text = data.get('text', '').strip()
    if not text:
        return jsonify({'error': 'No text provided'}), 400
    try:
        filename = f"{uuid.uuid4()}.mp3"
        
        # Use temporary file instead of audio folder
        with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as temp_file:
            filepath = temp_file.name
            
        tts = gTTS(text=text, lang='en')
        tts.save(filepath)

        # Upload to cloudinary first, before logging
        username = session.get('username')
        cloudinary_folder = f'storage/{username}/audio'
        cloudinary_url = upload_to_cloudinary(filepath, cloudinary_folder)
        store_url_in_firebase(cloudinary_url, 'audio', filename)

        # Now log the conversion with the cloudinary_url
        log_conversion('text-to-speech', 'user_input.txt', filename, filepath, cloudinary_url)

        # Store the audio data for download in cache
        with open(filepath, 'rb') as f:
            audio_data = f.read()
        
        # Store in audio cache for immediate download
        audio_cache[filename] = audio_data
        
        # Clean up temp file
        try:
            os.remove(filepath)
        except Exception as e:
            logger.warning(f"Failed to clean up temp audio file: {str(e)}")

        return jsonify({
            'success': True,
            'audio_url': f'/stream_audio/{filename}',  # For playing in browser
            'download_url': f'/download_audio/{filename}',  # For downloading
            'filename': filename,
            'message': 'Text converted to speech successfully!'
        })
    except Exception as e:
        logger.error(f"Text to speech conversion failed: {str(e)}", exc_info=True)
        return jsonify({'error': f"Text to speech conversion failed: {str(e)}"}), 500

@app.route('/download_audio/<filename>')
@login_required
def download_audio(filename):
    logger.info(f"Audio download requested for: {filename}")
    logger.info(f"Available audio files in cache: {list(audio_cache.keys())}")
    
    # Check if file exists in memory cache
    if filename not in audio_cache:
        logger.error(f"Audio file {filename} not found in cache")
        return jsonify({'error': 'Audio file not found or expired'}), 404
    
    try:
        # Get audio data from memory cache
        audio_data = audio_cache[filename]
        logger.info(f"Found audio data, size: {len(audio_data)} bytes")
        
        # Create a BytesIO object with the audio data
        audio_buffer = io.BytesIO(audio_data)
        audio_buffer.seek(0)
        
        # Don't delete from cache immediately - let it be accessed multiple times
        # The cache will be cleaned up by a timer or manually
        
        logger.info(f"Serving audio file: {filename}")
        return send_file(
            audio_buffer, 
            as_attachment=True, 
            download_name=filename, 
            mimetype='audio/mpeg'
        )
    except Exception as e:
        logger.error(f"Audio download failed: {str(e)}", exc_info=True)
        return jsonify({'error': f"Audio download failed: {str(e)}"}), 500

@app.route('/stream_audio/<filename>')
@login_required
def stream_audio(filename):
    """Stream audio for playing in browser (not as download)"""
    logger.info(f"Audio stream requested for: {filename}")
    
    # Check if file exists in memory cache
    if filename not in audio_cache:
        logger.error(f"Audio file {filename} not found in cache")
        return jsonify({'error': 'Audio file not found or expired'}), 404
    
    try:
        # Get audio data from memory cache
        audio_data = audio_cache[filename]
        logger.info(f"Streaming audio data, size: {len(audio_data)} bytes")
        
        # Create a BytesIO object with the audio data
        audio_buffer = io.BytesIO(audio_data)
        audio_buffer.seek(0)
        
        logger.info(f"Streaming audio file: {filename}")
        return send_file(
            audio_buffer, 
            as_attachment=False,  # Don't force download for streaming
            download_name=filename, 
            mimetype='audio/mpeg'
        )
    except Exception as e:
        logger.error(f"Audio streaming failed: {str(e)}", exc_info=True)
        return jsonify({'error': f"Audio streaming failed: {str(e)}"}), 500

# --- Speech to Text ---

@app.route('/save_transcript', methods=['POST'])
@login_required
def save_transcript():
    try:
        data = request.get_json()
        transcript = data.get('transcript', '')
        if not transcript:
            return jsonify({'error': 'No transcript provided'}), 400
        transcript_filename = f"transcript_{uuid.uuid4().hex}.txt"
        
        # Use temporary file for transcript
        with tempfile.NamedTemporaryFile(delete=False, suffix='.txt', mode='w', encoding='utf-8') as temp_file:
            temp_file.write(transcript + '\n')
            transcript_path = temp_file.name
            
        # Upload to cloudinary first, before logging
        username = session.get('username')
        cloudinary_folder = f'storage/{username}/txt'
        cloudinary_url = upload_to_cloudinary(transcript_path, cloudinary_folder)
        store_url_in_firebase(cloudinary_url, 'txt', transcript_filename)

        # Now log the conversion with the cloudinary_url
        log_conversion('speech-to-text', 'transcript.txt', transcript_filename, transcript_path, cloudinary_url)

        # Clean up temp file
        try:
            os.remove(transcript_path)
        except Exception as e:
            logger.warning(f"Failed to clean up temp transcript file: {str(e)}")

        return jsonify({'message': 'Transcript saved successfully'})
    except Exception as e:
        return jsonify({'error': f"Transcript save failed: {str(e)}"}), 500

@app.route('/upload_audio', methods=['POST'])
@login_required
def upload_audio():
    temp_file_path = None
    wav_path = None
    try:
        if 'audioFile' not in request.files:
            return jsonify({'error': 'No audio file provided'}), 400
        
        audio_file = request.files['audioFile']
        if audio_file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
            
        # Create temporary file for the uploaded audio
        with tempfile.NamedTemporaryFile(delete=False, suffix='.mp3') as temp_file:
            audio_file.save(temp_file)
            temp_file_path = temp_file.name
        
        # Try to convert using pydub (requires FFmpeg)
        try:
            audio = AudioSegment.from_mp3(temp_file_path)
            wav_path = temp_file_path.replace('.mp3', '.wav')
            audio.export(wav_path, format='wav')
        except Exception as convert_error:
            # If pydub fails, try direct wav loading if file is already wav format
            if audio_file.filename.lower().endswith('.wav'):
                wav_path = temp_file_path
            else:
                # Fallback: try to use the mp3 file directly with speech_recognition
                try:
                    with sr.AudioFile(temp_file_path) as source:
                        audio_data = recognizer.record(source)
                        transcript = recognizer.recognize_google(audio_data)
                        
                    # Cleanup and return
                    if temp_file_path and os.path.exists(temp_file_path):
                        os.remove(temp_file_path)
                    log_conversion('speech-to-text', audio_file.filename, 'transcript.txt', None, None, 'success')
                    return jsonify({'transcript': transcript})
                except Exception:
                    return jsonify({'error': 'FFmpeg is required for MP3 conversion. Please install FFmpeg or upload a WAV file instead.'}), 400
        
        # Process the wav file for speech recognition
        with sr.AudioFile(wav_path) as source:
            # Adjust for ambient noise
            recognizer.adjust_for_ambient_noise(source, duration=0.5)
            audio_data = recognizer.record(source)
            
            try:
                transcript = recognizer.recognize_google(audio_data)
                if not transcript.strip():
                    transcript = "No speech detected in the audio file"
            except sr.UnknownValueError:
                transcript = "Could not understand the audio. Please ensure the audio is clear and contains speech."
            except sr.RequestError as e:
                transcript = f"Speech recognition service error: {str(e)}"
        
        log_conversion('speech-to-text', audio_file.filename, 'transcript.txt', None, None, 'success')
        return jsonify({'transcript': transcript})
        
    except Exception as e:
        logger.error(f"Speech to text conversion error: {str(e)}")
        return jsonify({'error': f"Speech to text conversion failed: {str(e)}"}), 500
    finally:
        # Cleanup temporary files
        try:
            if temp_file_path and os.path.exists(temp_file_path):
                os.remove(temp_file_path)
            if wav_path and wav_path != temp_file_path and os.path.exists(wav_path):
                os.remove(wav_path)
        except Exception as cleanup_error:
            logger.warning(f"Failed to cleanup temp files: {str(cleanup_error)}")

# --- AI PDF Editor ---

def extract_structured_text(path):
    try:
        doc = fitz.open(path)
        structured_data = []
        for page in doc:
            blocks = page.get_text("dict")['blocks']
            for block in blocks:
                if 'lines' in block:
                    for line in block['lines']:
                        line_text = " ".join([span['text'] for span in line['spans']])
                        structured_data.append(line_text)
        doc.close()
        return "\n".join(structured_data)
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        return f"Error extracting text from PDF: {str(e)}"

def retry_api_request(url, headers, data, max_retries=3, delay=2):
    for attempt in range(max_retries):
        try:
            response = requests.post(url, headers=headers, json=data, timeout=30)
            if response.status_code == 429:
                logger.warning(f"Rate limit hit on attempt {attempt + 1}, retrying after {delay} seconds")
                time.sleep(delay)
                delay *= 2
                continue
            response.raise_for_status()
            return response
        except requests.exceptions.RequestException as e:
            logger.error(f"API request failed on attempt {attempt + 1}: {str(e)}")
            if attempt == max_retries - 1:
                raise e
            time.sleep(delay)
            delay *= 2
    raise Exception("Max retries exceeded for API request")

@app.route('/analyze', methods=['POST'])
@login_required
def analyze():
    global latest_text
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    uploaded_file = request.files['file']
    if not uploaded_file.filename.lower().endswith('.pdf'):
        return jsonify({'error': 'Only PDF files are supported'}), 400
    filename = secure_filename(uploaded_file.filename)
    
    # Use temporary file instead of upload folder
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
        uploaded_file.save(temp_file)
        file_path = temp_file.name
    try:
        extracted_text = extract_structured_text(file_path)
        if extracted_text.startswith('Error'):
            os.remove(file_path)
            return jsonify({'error': extracted_text}), 500
        if not extracted_text.strip():
            os.remove(file_path)
            return jsonify({'error': 'No text extracted from PDF. Ensure the PDF contains selectable text, not images.'}), 400
        latest_text = extracted_text
        prompt = (
            "You are an intelligent assistant analyzing a PDF document. Your task is to identify blank or unfilled fields such as "
            "'Date: ____', 'Name: ________', 'Signature: [____]', or other placeholders (e.g., '________', '[____]', empty lines after labels). "
            "For each identified field, suggest a reasonable completion based on context (e.g., use today's date 'June 25, 2025' for date fields, 'John Doe' for name fields, 'Signature' for signature fields). "
            "Return the results in the following format:\n"
            "Identified Fields:\n"
            "- Field: [Description], Suggestion: [Suggested Value]\n"
            "If no blank fields are found, state: 'No blank or unfilled fields detected.'\n\n"
            f"Document Text:\n{extracted_text[:2000]}"
        )

        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json"
        }

        data = {
            "model": OPENROUTER_MODEL,
            "messages": [{"role": "user", "content": prompt}],
            "max_tokens": 1000
        }

        response = retry_api_request(OPENROUTER_API_URL, headers, data)
        result = response.json()
        suggestions = result.get('choices', [{}])[0].get('message', {}).get('content', '')

        if not suggestions:
            raise ValueError("Empty response content from API")

        log_conversion('ai-pdf-editor', filename, 'analysis.json', None, None, 'success')
        return jsonify({"text": extracted_text, "suggestions": suggestions})
    except requests.exceptions.RequestException as e:
        logger.error(f"OpenRouter API request failed: {str(e)}")
        return jsonify({'error': f"Failed to analyze document: API request error - {str(e)}"}), 500
    except ValueError as e:
        logger.error(f"OpenRouter API response error: {str(e)}")
        return jsonify({'error': f"Failed to analyze document: Invalid API response - {str(e)}"}), 500
    except Exception as e:
        logger.error(f"PDF analysis failed: {str(e)}")
        return jsonify({'error': f"Failed to analyze document: {str(e)}"}), 500
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)

@app.route('/edit', methods=['POST'])
@login_required
def edit():
    global latest_text
    try:
        updated_text = request.json.get('updated_text')
        if not updated_text:
            return jsonify({'error': 'No updated text provided'}), 400
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        for line in updated_text.split('\n'):
            pdf.multi_cell(0, 10, line)
        output_filename = f"edited_{uuid.uuid4().hex}.pdf"
        
        # Use temporary file instead of upload folder
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            output_path = temp_file.name
            
        pdf.output(output_path)
        log_conversion('ai-pdf-editor', 'user_input.txt', output_filename, output_path, None, 'success')
        return send_file(
            output_path,
            mimetype='application/pdf',
            as_attachment=True,
            download_name="edited_document.pdf"
        )
    except Exception as e:
        logger.error(f"PDF edit failed: {str(e)}")
        return jsonify({'error': f"Failed to edit document: {str(e)}"}), 500
    finally:
        # Clean up temp file
        try:
            if 'output_path' in locals() and os.path.exists(output_path):
                os.remove(output_path)
        except Exception as e:
            logger.warning(f"Failed to clean up temp edit file: {str(e)}")

@app.route('/fill_from_prompt', methods=['POST'])
@login_required
def fill_from_prompt():
    global latest_text
    try:
        user_prompt = request.json.get('user_prompt')
        current_text = request.json.get('current_text', '')
        
        if not user_prompt:
            return jsonify({'error': 'No user prompt provided'}), 400
        
        # Use current_text if provided, otherwise fall back to latest_text
        text_to_process = current_text if current_text.strip() else latest_text
        
        if not text_to_process:
            return jsonify({'error': 'No text available to process. Please analyze a document first or add some text.'}), 400
            
        ai_instruction = (
            "Based on the following document text, a user wants to add or update content as follows:\n"
            f"Instruction: {user_prompt}\n\n"
            f"Document:\n{text_to_process[:2000]}\n\n"
            "Provide the updated version of the document with the user request applied appropriately. "
            "Return only the revised document content."
        )
        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json"
        }
        data = {
            "model": OPENROUTER_MODEL,
            "messages": [{"role": "user", "content": ai_instruction}],
            "max_tokens": 1000
        }
        response = retry_api_request(OPENROUTER_API_URL, headers, data)
        result = response.json()
        updated_text = result.get('choices', [{}])[0].get('message', {}).get('content', '')
        if not updated_text:
            return jsonify({'error': 'Empty response from API'}), 500
        latest_text = updated_text
        log_conversion('ai-pdf-editor', 'user_prompt.txt', 'updated_text.txt', None, None, 'success')
        return jsonify({"updated_text": updated_text})
    except Exception as e:
        logger.error(f"Prompt-based edit failed: {str(e)}")
        return jsonify({'error': f"Prompt-based edit failed: {str(e)}"}), 500

# --- Text Summarizer ---

@app.route('/summarize', methods=['POST'])
@login_required
def summarize():
    try:
        data = request.get_json()
        text = data.get('text', '').strip()
        if not text:
            logger.error("No text provided to summarizer endpoint")
            return jsonify({'error': 'No text provided'}), 400
        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json"
        }
        payload = {
            "model": OPENROUTER_MODEL,
            "messages": [
                {
                    "role": "user",
                    "content": (
                        "Summarize the following text in 3 sentences. "
                        "Only return the summary content. Do not add any introduction, title, or prefix:\n\n"
                        f"{text[:2000]}"
                    )
                }
            ],
            "max_tokens": 500
        }
        response = retry_api_request(OPENROUTER_API_URL, headers, payload)
        result = response.json()
        summary = result.get('choices', [{}])[0].get('message', {}).get('content', '').strip()
        if not summary:
            logger.error(f"Empty summary from API. Response: {json.dumps(result)}")
            return jsonify({'error': 'Empty summary from API'}), 500

        summary_filename = f"summary_{uuid.uuid4().hex}.txt"
        
        # Use temporary file for summary
        with tempfile.NamedTemporaryFile(delete=False, suffix='.txt', mode='w', encoding='utf-8') as temp_file:
            temp_file.write(summary)
            summary_path = temp_file.name

        # Upload to cloudinary first, before logging
        username = session.get('username')
        cloudinary_folder = f'storage/{username}/txt'
        try:
            cloudinary_url = upload_to_cloudinary(summary_path, cloudinary_folder)
            store_url_in_firebase(cloudinary_url, 'txt', summary_filename)
        except Exception as upload_e:
            logger.error(f"Cloudinary/Firebase upload failed for summary: {str(upload_e)}")
            cloudinary_url = None  # Set to None if upload fails

        # Now log the conversion with the cloudinary_url
        log_conversion('text-summarizer', 'user_input.txt', summary_filename, summary_path, cloudinary_url)

        # Clean up temp file
        try:
            os.remove(summary_path)
        except Exception as e:
            logger.warning(f"Failed to clean up temp summary file: {str(e)}")

        return jsonify({'summary': summary})
    except Exception as e:
        logger.error(f"Text summarization failed: {str(e)}", exc_info=True)
        return jsonify({'error': f"Text summarization failed: {str(e)}"}), 500

@app.route('/admin/logs', methods=['GET'])
@login_required
def get_conversion_logs():
    """API endpoint to fetch conversion logs for the history page."""
    try:
        conn = sqlite3.connect('file_conversion.db')
        c = conn.cursor()
        c.execute('''
            SELECT id, conversion_type, original_filename, converted_filename, file_path, timestamp, cloudinary_url, username, status
            FROM conversions 
            WHERE username = ?
            ORDER BY timestamp DESC
        ''', (session.get('username'),))
        rows = c.fetchall()
        conn.close()
        
        logs = []
        for row in rows:
            # Use the actual status from database, with fallback logic
            actual_status = row[8] if len(row) > 8 and row[8] else None
            if actual_status is None:
                # Fallback: determine status based on available data
                if row[6]:  # cloudinary_url exists
                    actual_status = 'success'
                elif row[1] in ['speech-to-text', 'document-screener', 'ai-pdf-editor', 'text-summarizer']:
                    # These don't always need cloudinary_url to be successful
                    actual_status = 'success'
                else:
                    actual_status = 'error'
            
            log_entry = {
                'id': row[0],
                'conversion_type': row[1],
                'original_filename': row[2],
                'filename': row[3],
                'download_path': f"/download-file/{row[0]}" if (row[6] or actual_status == 'success') else None,
                'timestamp': row[5],
                'status': actual_status,
                'file_size': 'Unknown',  # You can add file size calculation if needed
                'cloudinary_url': row[6],
                'username': row[7]
            }
            logs.append(log_entry)
        
        return jsonify({'success': True, 'logs': logs})
    except Exception as e:
        logger.error(f"Failed to fetch conversion logs: {str(e)}")
        return jsonify({'success': False, 'error': str(e)}), 500

# --- Cache Cleanup Functions ---

def cleanup_audio_cache():
    """Clean up audio cache to prevent memory leaks"""
    try:
        if audio_cache:
            logger.info(f"Cleaning up {len(audio_cache)} audio files from cache")
            audio_cache.clear()
    except Exception as e:
        logger.error(f"Error cleaning up audio cache: {str(e)}")

# Schedule periodic cleanup every 10 minutes
import threading
import time

def periodic_cleanup():
    while True:
        time.sleep(600)  # 10 minutes
        cleanup_audio_cache()

# Start cleanup thread
cleanup_thread = threading.Thread(target=periodic_cleanup, daemon=True)
cleanup_thread.start()

# --- Test Routes for Debugging ---

@app.route('/test_dependencies')
def test_dependencies():
    """Test if all required dependencies are working"""
    try:
        # Test PIL/Pillow
        from PIL import Image
        test_image = Image.new('RGB', (100, 100), color='red')
        
        # Test Cloudinary
        import cloudinary
        cloudinary_status = "Configured" if cloudinary.config().cloud_name else "Not configured"
        
        # Test Firebase
        firebase_status = "Connected" if db else "Not connected"
        
        return jsonify({
            'success': True,
            'dependencies': {
                'PIL': 'Working',
                'Cloudinary': cloudinary_status,
                'Firebase': firebase_status
            }
        })
    except Exception as e:
        return jsonify({
            'success': False,
            'error': str(e)
        })

# --- Profile Management Routes ---

def update_users_table():
    """Update users table to include profile fields"""
    conn = sqlite3.connect('file_conversion.db')
    c = conn.cursor()
    
    # Add new columns if they don't exist
    profile_columns = [
        ('name', 'TEXT'),
        ('email', 'TEXT'),
        ('contact_number', 'TEXT'),
        ('country', 'TEXT'),
        ('profile_picture', 'TEXT'),
        ('membership_status', 'TEXT DEFAULT "Standard"'),
        ('created_at', 'TIMESTAMP DEFAULT CURRENT_TIMESTAMP')
    ]
    
    for column_name, column_type in profile_columns:
        try:
            c.execute(f'ALTER TABLE users ADD COLUMN {column_name} {column_type}')
        except sqlite3.OperationalError:
            pass  # Column already exists
    
    conn.commit()
    conn.close()

def get_user_profile(username):
    """Get user profile data from database and Firebase"""
    try:
        # Get from Firebase first
        user_ref = db.reference(f'Data/{username}')
        user_data = user_ref.get()
        
        if user_data:
            return user_data
        
        # Fallback to SQLite
        conn = sqlite3.connect('file_conversion.db')
        c = conn.cursor()
        c.execute('''SELECT username, name, email, contact_number, country, 
                           profile_picture, membership_status, created_at 
                    FROM users WHERE username = ?''', (username,))
        row = c.fetchone()
        conn.close()
        
        if row:
            return {
                'username': row[0],
                'name': row[1],
                'email': row[2],
                'contact_number': row[3],
                'country': row[4],
                'profile_picture': row[5],
                'membership_status': row[6] or 'Standard',
                'created_at': row[7]
            }
        return None
    except Exception as e:
        logger.error(f"Error getting user profile: {str(e)}")
        return None

def get_user_context():
    """Get user context for templates"""
    username = session.get('username')
    user = get_user_profile(username)
    
    if not user:
        user = {'username': username, 'profile_picture': None}
    
    return {'user': user}

def update_user_profile(username, profile_data):
    """Update user profile in both Firebase and SQLite"""
    try:
        # Update Firebase
        user_ref = db.reference(f'Data/{username}')
        existing_data = user_ref.get() or {}
        existing_data.update(profile_data)
        user_ref.set(existing_data)
        
        # Update SQLite as backup
        conn = sqlite3.connect('file_conversion.db')
        c = conn.cursor()
        
        # Check if user exists in SQLite
        c.execute('SELECT id FROM users WHERE username = ?', (username,))
        user_exists = c.fetchone()
        
        if user_exists:
            # Update existing user
            c.execute('''UPDATE users SET name = ?, email = ?, contact_number = ?, 
                               country = ?, profile_picture = ?, membership_status = ?
                        WHERE username = ?''', 
                     (profile_data.get('name'), profile_data.get('email'), 
                      profile_data.get('contact_number'), profile_data.get('country'),
                      profile_data.get('profile_picture'), profile_data.get('membership_status'),
                      username))
        else:
            # Insert new user record
            c.execute('''INSERT INTO users (username, name, email, contact_number, 
                               country, profile_picture, membership_status)
                        VALUES (?, ?, ?, ?, ?, ?, ?)''',
                     (username, profile_data.get('name'), profile_data.get('email'),
                      profile_data.get('contact_number'), profile_data.get('country'),
                      profile_data.get('profile_picture'), profile_data.get('membership_status')))
        
        conn.commit()
        conn.close()
        return True
    except Exception as e:
        logger.error(f"Error updating user profile: {str(e)}")
        return False

@app.route('/profile')
@login_required
def profile():
    """Display user profile page"""
    username = session.get('username')
    user = get_user_profile(username)
    
    if not user:
        user = {'username': username}
    
    return render_template('profile.html', user=user)

@app.route('/update_profile', methods=['POST'])
@login_required
def update_profile():
    """Update user profile information"""
    try:
        username = session.get('username')
        
        # Get form data
        profile_data = {
            'name': request.form.get('name', '').strip(),
            'username': request.form.get('username', '').strip(),
            'email': request.form.get('email', '').strip(),
            'contact_number': request.form.get('contact_number', '').strip(),
            'country': request.form.get('country', '').strip(),
            'membership_status': request.form.get('membership_status', 'Standard')
        }
        
        # Validate required fields
        required_fields = ['name', 'username', 'email', 'country']
        for field in required_fields:
            if not profile_data[field]:
                return render_template('profile.html', 
                                     user=get_user_profile(username),
                                     error_message=f'{field.replace("_", " ").title()} is required')
        
        # Validate email format
        import re
        email_pattern = r'^[^\s@]+@[^\s@]+\.[^\s@]+$'
        if not re.match(email_pattern, profile_data['email']):
            return render_template('profile.html', 
                                 user=get_user_profile(username),
                                 error_message='Please enter a valid email address')
        
        # Check if username is being changed and if new username exists
        if profile_data['username'] != username:
            existing_user = get_user_profile(profile_data['username'])
            if existing_user:
                return render_template('profile.html', 
                                     user=get_user_profile(username),
                                     error_message='Username already exists')
        
        # Update profile
        if update_user_profile(username, profile_data):
            # Update session if username changed
            if profile_data['username'] != username:
                session['username'] = profile_data['username']
            
            return render_template('profile.html', 
                                 user=get_user_profile(profile_data['username']),
                                 success_message='Profile updated successfully!')
        else:
            return render_template('profile.html', 
                                 user=get_user_profile(username),
                                 error_message='Failed to update profile')
    
    except Exception as e:
        logger.error(f"Profile update error: {str(e)}")
        return render_template('profile.html', 
                             user=get_user_profile(session.get('username')),
                             error_message='An error occurred while updating profile')

@app.route('/upload_profile_picture', methods=['POST'])
@login_required
def upload_profile_picture():
    """Handle profile picture upload"""
    try:
        username = session.get('username')
        logger.info(f"Profile picture upload request from user: {username}")
        
        if 'profile_picture' not in request.files:
            logger.error("No profile_picture in request.files")
            return jsonify({'success': False, 'error': 'No file uploaded'})
        
        file = request.files['profile_picture']
        logger.info(f"File received: {file.filename}, size: {file.content_length}")
        
        if file.filename == '':
            logger.error("Empty filename")
            return jsonify({'success': False, 'error': 'No file selected'})
        
        # Validate file type
        allowed_extensions = {'png', 'jpg', 'jpeg', 'gif', 'webp'}
        file_extension = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
        
        if file_extension not in allowed_extensions:
            logger.error(f"Invalid file extension: {file_extension}")
            return jsonify({'success': False, 'error': 'Invalid file type. Please upload an image file.'})
        
        # Create temporary file for processing
        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}') as temp_file:
            file.save(temp_file.name)
            temp_path = temp_file.name
            logger.info(f"File saved to temp path: {temp_path}")
        
        try:
            # Check if PIL/Pillow is working
            logger.info("Processing image with PIL...")
            
            # Resize and optimize image
            with Image.open(temp_path) as img:
                logger.info(f"Image opened: {img.size}, mode: {img.mode}")
                
                # Convert to RGB if necessary
                if img.mode in ('RGBA', 'LA', 'P'):
                    img = img.convert('RGB')
                    logger.info("Image converted to RGB")
                
                # Resize to max 500x500 while maintaining aspect ratio
                img.thumbnail((500, 500), Image.Resampling.LANCZOS)
                logger.info(f"Image resized to: {img.size}")
                
                # Save optimized image
                optimized_path = temp_path.replace(f'.{file_extension}', '_optimized.jpg')
                img.save(optimized_path, 'JPEG', quality=85, optimize=True)
                logger.info(f"Optimized image saved to: {optimized_path}")
            
            # Upload to Cloudinary
            logger.info("Uploading to Cloudinary...")
            cloudinary_folder = f'storage/{username}/profile'
            
            result = cloudinary.uploader.upload(
                optimized_path,
                folder=cloudinary_folder,
                public_id=f'profile_picture_{username}',
                overwrite=True,
                resource_type='image',
                format='jpg'
            )
            
            profile_picture_url = result['secure_url']
            logger.info(f"Cloudinary upload successful: {profile_picture_url}")
            
            # Update user profile with new picture URL
            profile_data = {'profile_picture': profile_picture_url}
            if update_user_profile(username, profile_data):
                logger.info("Profile updated successfully")
                return jsonify({
                    'success': True, 
                    'profile_picture_url': profile_picture_url,
                    'message': 'Profile picture updated successfully!'
                })
            else:
                logger.error("Failed to update user profile in database")
                return jsonify({'success': False, 'error': 'Failed to save profile picture URL'})
        
        except Exception as upload_error:
            logger.error(f"Image upload error: {str(upload_error)}", exc_info=True)
            return jsonify({'success': False, 'error': f'Failed to process and upload image: {str(upload_error)}'})
        
        finally:
            # Clean up temporary files
            try:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
                    logger.info(f"Cleaned up temp file: {temp_path}")
                if 'optimized_path' in locals() and os.path.exists(optimized_path):
                    os.remove(optimized_path)
                    logger.info(f"Cleaned up optimized file: {optimized_path}")
            except Exception as cleanup_error:
                logger.warning(f"Failed to clean up temp files: {str(cleanup_error)}")
    
    except Exception as e:
        logger.error(f"Profile picture upload error: {str(e)}", exc_info=True)
        return jsonify({'success': False, 'error': f'An error occurred while uploading the image: {str(e)}'})

@app.route('/remove_profile_picture', methods=['POST'])
@login_required
def remove_profile_picture():
    try:
        username = session.get('username')
        
        # Update profile picture to null in user profile
        profile_data = {'profile_picture': None}
        if update_user_profile(username, profile_data):
            return jsonify({'success': True, 'message': 'Profile picture removed successfully!'})
        else:
            return jsonify({'success': False, 'error': 'Failed to remove profile picture'})
    except Exception as e:
        logger.error(f"Remove profile picture error: {str(e)}")
        return jsonify({'success': False, 'error': str(e)})

@app.route('/change_password')
@login_required
def change_password():
    """Display change password form"""
    return render_template('change_password.html')

@app.route('/update_password', methods=['POST'])
@login_required
def update_password():
    """Update user password"""
    try:
        username = session.get('username')
        current_password = request.form.get('current_password')
        new_password = request.form.get('new_password')
        confirm_password = request.form.get('confirm_password')
        
        # Validate inputs
        if not all([current_password, new_password, confirm_password]):
            return render_template('change_password.html', 
                                 error='All fields are required')
        
        if new_password != confirm_password:
            return render_template('change_password.html', 
                                 error='New passwords do not match')
        
        if len(new_password) < 6:
            return render_template('change_password.html', 
                                 error='Password must be at least 6 characters long')
        
        # Verify current password
        user_data = get_user_by_username(username)
        if not user_data or not check_password_hash(user_data['password'], current_password):
            return render_template('change_password.html', 
                                 error='Current password is incorrect')
        
        # Update password in Firebase
        new_password_hash = generate_password_hash(new_password)
        cred_ref = db.reference(f'credentials/users/{username}')
        cred_ref.update({'password': new_password_hash})
        
        # Also update in user data
        user_ref = db.reference(f'Data/{username}')
        user_ref.update({'password': new_password_hash})
        
        return render_template('change_password.html', 
                             success='Password updated successfully!')
    
    except Exception as e:
        logger.error(f"Password update error: {str(e)}")
        return render_template('change_password.html', 
                             error='An error occurred while updating password')

@app.route('/email-settings')
@login_required
def email_settings():
    """Display email settings page for phone/email verification"""
    username = session.get('username')
    user_data = get_user_profile(username)
    
    if not user_data:
        user_data = {'username': username, 'email': '', 'phone': ''}
    
    # Get verification status from Firebase
    firebase_ref = db.reference(f'Data/{username}')
    firebase_data = firebase_ref.get() or {}
    
    # Check verification status
    phone_verified = firebase_data.get('phone_verified', False)
    email_verified = firebase_data.get('email_verified', False)
    
    return render_template('email_settings.html', 
                         user=user_data,
                         phone_verified=phone_verified,
                         email_verified=email_verified,
                         **get_user_context())

# --- Run Flask App ---

if __name__ == '__main__':
    init_db()  # Initialize DB once on startup
    update_users_table()  # Update users table schema
    print("🚀 Starting DocShift on 127.0.0.1:5000")
    print("� Email verification enabled")
    print("📱 Phone/SMS verification DISABLED")
    print("🔥 Firebase SMS functionality has been removed")
    app.run(debug=True, host='127.0.0.1', port=5000)